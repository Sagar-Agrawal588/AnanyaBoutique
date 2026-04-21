from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(r"d:\BogEcom")
PPT_PATH = ROOT / "Piyush_Songara_Mid_Review_Presentation.pptx"
OUTPUT_PATH = ROOT / "Piyush_Songara_Mid_Review_Presentation_GenZ_Theme.pptx"
LOGO = ROOT / "frontend" / "client" / "public" / "logo.png"
SCREEN_DIR = ROOT / "output" / "mid-review-live-screenshots"
JECRC_LOGO = ROOT / "output" / "jecrc-assets" / "JU-Logo.png"
JECRC_FOUNDATION = ROOT / "output" / "jecrc-assets" / "jecrc-foundation.png"

BG = RGBColor(245, 241, 232)
CARD = RGBColor(255, 252, 247)
ACCENT = RGBColor(237, 225, 188)
ACCENT_DARK = RGBColor(182, 147, 66)
TEXT = RGBColor(78, 61, 46)
MUTED = RGBColor(106, 96, 84)
DARK_PANEL = RGBColor(68, 62, 63)
WHITE = RGBColor(255, 255, 255)
BLUE_NOTE = RGBColor(213, 229, 252)
SOFT_LINE = RGBColor(226, 217, 201)
RED = RGBColor(238, 27, 36)
LIGHT_PINK = RGBColor(255, 245, 245)


def set_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_card(slide, left, top, width, height, fill=CARD, line=RGBColor(226, 217, 201), radius=True):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1.2)
    return shape


def add_text(slide, left, top, width, height, text, *, font="Georgia", size=20, color=TEXT, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Pt(2)
    tf.margin_right = Pt(2)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.space_after = Pt(0)
    p.line_spacing = 1.1
    for run in p.runs:
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return box


def add_multiline(slide, left, top, width, height, lines, *, font="Aptos", size=14, color=MUTED, bold_first=False, dark=False, align=PP_ALIGN.LEFT, gap=4):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Pt(2)
    tf.margin_right = Pt(2)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    tf.clear()
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        p.space_after = Pt(gap if idx < len(lines) - 1 else 0)
        p.line_spacing = 1.2
        for run in p.runs:
            run.font.name = font
            run.font.size = Pt(size)
            run.font.bold = bold_first and idx == 0
            run.font.color.rgb = WHITE if dark else color
    return box


def add_header_chrome(slide, section_no, label):
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.82), Inches(0.42), Inches(0.42), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_DARK
    line.line.width = Pt(0)
    chip = add_card(slide, Inches(11.35), Inches(0.42), Inches(1.15), Inches(0.32), fill=RGBColor(252, 247, 232), line=SOFT_LINE)
    p = chip.text_frame.paragraphs[0]
    p.text = f"{section_no}  {label}"
    p.alignment = PP_ALIGN.CENTER
    for run in p.runs:
        run.font.name = "Aptos"
        run.font.size = Pt(8.5)
        run.font.color.rgb = MUTED


def add_icon_badge(slide, left, top, text, dark=False):
    badge = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, left, top, Inches(0.28), Inches(0.28))
    badge.fill.solid()
    badge.fill.fore_color.rgb = ACCENT_DARK if not dark else RGBColor(104, 93, 93)
    badge.line.width = Pt(0)
    p = badge.text_frame.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    for run in p.runs:
        run.font.name = "Aptos"
        run.font.size = Pt(9)
        run.font.bold = True
        run.font.color.rgb = WHITE
    return badge


def add_jecrc_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    top = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.15))
    top.fill.solid()
    top.fill.fore_color.rgb = RED
    top.line.width = Pt(0)
    bottom = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(7.32), Inches(13.333), Inches(0.18))
    bottom.fill.solid()
    bottom.fill.fore_color.rgb = RED
    bottom.line.width = Pt(0)
    wash = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(1.15), Inches(13.333), Inches(6.17))
    wash.fill.solid()
    wash.fill.fore_color.rgb = LIGHT_PINK
    wash.fill.transparency = 0.55
    wash.line.width = Pt(0)


def add_watermark_words(slide):
    words = [
        ("TIME", 5.2, 2.05, 54),
        ("BUILD", 9.4, 4.15, 52),
        ("BROTHERHOOD", 4.3, 5.85, 42),
        ("THINK", 3.6, 4.0, 30),
        ("TEACHING", 4.5, 1.55, 26),
        ("INNOVATION", 10.6, 1.05, 24),
        ("INFORMATION", 7.8, 6.55, 28),
        ("KNOWLEDGE", 2.5, 4.6, 22),
        ("PROGRAM", 0.8, 2.8, 20),
        ("LIFE", 6.0, 1.8, 20),
    ]
    for text, x, y, size in words:
        add_text(slide, Inches(x), Inches(y), Inches(3.8), Inches(0.7), text, font="Georgia", size=size, color=RGBColor(223, 159, 159), align=PP_ALIGN.CENTER)
        slide.shapes[-1].text_frame.paragraphs[0].runs[0].font.color.transparency = 0.82


def add_jecrc_title_slide(slide):
    add_jecrc_background(slide)
    add_watermark_words(slide)
    add_text(slide, Inches(4.85), Inches(0.14), Inches(3.8), Inches(0.42), "DEPARTMENT OF", font="Aptos", size=26, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(5.2), Inches(0.58), Inches(3.2), Inches(0.35), "Computer Science", font="Aptos", size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    if JECRC_LOGO.exists():
        slide.shapes.add_picture(str(JECRC_LOGO), Inches(10.1), Inches(0.12), height=Inches(0.92))
        slide.shapes.add_picture(str(JECRC_LOGO), Inches(1.35), Inches(3.85), height=Inches(0.75))
    if LOGO.exists():
        slide.shapes.add_picture(str(LOGO), Inches(10.45), Inches(3.9), height=Inches(0.62))
    if JECRC_FOUNDATION.exists():
        slide.shapes.add_picture(str(JECRC_FOUNDATION), Inches(11.2), Inches(6.55), height=Inches(0.58))
    add_text(slide, Inches(3.95), Inches(1.48), Inches(5.2), Inches(0.5), "Mid Review Internship Presentation", font="Georgia", size=22, color=RGBColor(138, 138, 138), align=PP_ALIGN.CENTER)
    add_text(slide, Inches(5.45), Inches(2.1), Inches(2.3), Inches(0.3), "Presented by", font="Georgia", size=14, color=TEXT, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(5.1), Inches(2.38), Inches(3.0), Inches(0.38), "Piyush Songara", font="Georgia", size=20, color=RGBColor(0, 0, 0), bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(4.7), Inches(2.72), Inches(3.8), Inches(0.4), "Under the guidance of", font="Georgia", size=18, color=RGBColor(0, 0, 0), bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(0.25), Inches(5.02), Inches(5.2), Inches(0.42), "Faculty Internship guide:-Megha Garg", font="Georgia", size=18, color=RGBColor(0, 0, 0), bold=True, align=PP_ALIGN.LEFT)
    add_text(slide, Inches(8.6), Inches(5.02), Inches(4.3), Inches(0.42), "Industrial guide:-Balveer Bajiya", font="Georgia", size=18, color=RGBColor(0, 0, 0), bold=True, align=PP_ALIGN.LEFT)
    building = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(1.1), Inches(6.08), Inches(10.8), Inches(0.72))
    building.fill.solid()
    building.fill.fore_color.rgb = RGBColor(183, 117, 95)
    building.line.width = Pt(0)
    add_text(slide, Inches(5.35), Inches(7.02), Inches(2.5), Inches(0.25), "2025-2026", font="Aptos", size=16, color=RGBColor(0, 0, 0), align=PP_ALIGN.CENTER)


def add_jecrc_contents_slide(slide):
    add_jecrc_background(slide)
    add_watermark_words(slide)
    add_text(slide, Inches(5.25), Inches(0.22), Inches(2.8), Inches(0.35), "Contents", font="Georgia", size=22, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    if JECRC_LOGO.exists():
        slide.shapes.add_picture(str(JECRC_LOGO), Inches(10.1), Inches(0.12), height=Inches(0.92))
    if JECRC_FOUNDATION.exists():
        slide.shapes.add_picture(str(JECRC_FOUNDATION), Inches(11.2), Inches(6.55), height=Inches(0.58))
    items = [
        "Company / Project Profile",
        "Internship Role",
        "Project Introduction",
        "Work Evidence and Metrics",
        "Roles and Responsibilities",
        "Tools Used",
        "Methodology",
        "Outcomes and Learning",
        "Conclusion",
    ]
    y = 1.55
    for item in items:
        add_text(slide, Inches(0.28), Inches(y), Inches(6.0), Inches(0.35), item, font="Georgia", size=18, color=RGBColor(0, 0, 0), align=PP_ALIGN.LEFT)
        y += 0.6


def add_tag(slide, left, top, text):
    tag = add_card(slide, left, top, Inches(1.55), Inches(0.34), fill=RGBColor(252, 247, 232), line=ACCENT_DARK)
    tf = tag.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    for run in p.runs:
        run.font.name = "Aptos"
        run.font.size = Pt(9.5)
        run.font.color.rgb = ACCENT_DARK


def add_split_hero(slide, image_path, eyebrow, title, subtitle, meta, tags):
    set_bg(slide)
    add_card(slide, Inches(0.5), Inches(0.45), Inches(12.35), Inches(6.2))
    add_header_chrome(slide, "01", "Cover")
    left_panel = add_card(slide, Inches(0.5), Inches(0.45), Inches(4.6), Inches(6.2), fill=RGBColor(248, 237, 202), line=RGBColor(238, 226, 191), radius=False)
    left_panel.line.width = Pt(0)
    if image_path and image_path.exists():
        slide.shapes.add_picture(str(image_path), Inches(1.0), Inches(1.15), width=Inches(3.55), height=Inches(3.9))
    add_text(slide, Inches(5.75), Inches(0.95), Inches(5.75), Inches(0.5), eyebrow, font="Georgia", size=18, color=MUTED)
    add_text(slide, Inches(5.75), Inches(1.45), Inches(5.9), Inches(1.8), title, font="Georgia", size=34, color=TEXT, bold=False)
    add_text(slide, Inches(5.75), Inches(3.55), Inches(5.9), Inches(0.55), subtitle, font="Aptos", size=16, color=TEXT, bold=True)
    add_text(slide, Inches(5.75), Inches(4.03), Inches(6.0), Inches(0.55), meta, font="Aptos", size=14, color=MUTED)
    x = Inches(5.75)
    for tag in tags:
        add_tag(slide, x, Inches(4.65), tag)
        x += Inches(1.72)
    if LOGO.exists():
        slide.shapes.add_picture(str(LOGO), Inches(11.6), Inches(0.72), width=Inches(0.7), height=Inches(0.7))


def add_overview_slide(slide, image_path):
    set_bg(slide)
    add_card(slide, Inches(0.55), Inches(0.45), Inches(12.2), Inches(6.15))
    add_header_chrome(slide, "02", "Overview")
    accent = add_card(slide, Inches(0.55), Inches(0.45), Inches(4.6), Inches(6.15), fill=RGBColor(248, 237, 202), line=RGBColor(238, 226, 191), radius=False)
    accent.line.width = Pt(0)
    if image_path.exists():
        slide.shapes.add_picture(str(image_path), Inches(0.95), Inches(1.55), width=Inches(3.7), height=Inches(3.2))
    add_text(slide, Inches(5.75), Inches(0.75), Inches(5.6), Inches(0.7), "Abstract & Overview", font="Georgia", size=26, color=TEXT)
    add_text(
        slide,
        Inches(5.75),
        Inches(1.55),
        Inches(5.85),
        Inches(1.05),
        "This presentation documents a structured internship on BogEcom (Buy One Gram), a live ecommerce platform. My role covered end-to-end UI implementation, state management, production fixes, API-linked behavior, and delivery support across client and admin portals.",
        font="Aptos",
        size=15,
        color=MUTED,
    )
    cards = [
        ("Role", "Frontend Developer Intern with full product exposure from day one"),
        ("Platform", "Dual-portal ecommerce system with client storefront, admin dashboard, and connected APIs"),
        ("Key Contributions", "UI components, checkout flows, state management, API integration, bug fixes"),
        ("Duration", "January 2026 - April 2026, including latest updates up to April 21"),
    ]
    top = 2.55
    for idx, (title, body) in enumerate(cards, start=1):
        add_card(slide, Inches(5.75), Inches(top), Inches(5.7), Inches(0.9), fill=RGBColor(249, 239, 210), line=RGBColor(228, 208, 150))
        add_icon_badge(slide, Inches(5.92), Inches(top + 0.16), str(idx))
        add_text(slide, Inches(5.95), Inches(top + 0.12), Inches(2.0), Inches(0.3), title, font="Georgia", size=16, color=TEXT)
        add_text(slide, Inches(5.95), Inches(top + 0.42), Inches(5.2), Inches(0.3), body, font="Aptos", size=11.5, color=MUTED)
        top += 1.02


def add_dual_panel(slide, title, subtitle, left_title, left_lines, right_title, right_lines):
    set_bg(slide)
    add_header_chrome(slide, "03", "Context")
    add_text(slide, Inches(0.8), Inches(0.58), Inches(6.0), Inches(0.55), title, font="Georgia", size=24, color=TEXT)
    add_text(slide, Inches(0.8), Inches(1.12), Inches(7.5), Inches(0.4), subtitle, font="Aptos", size=12.5, color=MUTED)
    add_card(slide, Inches(0.8), Inches(1.8), Inches(4.15), Inches(3.85), fill=DARK_PANEL, line=DARK_PANEL)
    add_text(slide, Inches(1.05), Inches(2.05), Inches(3.5), Inches(0.45), left_title, font="Georgia", size=18, color=WHITE, bold=True)
    add_multiline(slide, Inches(1.05), Inches(2.52), Inches(3.45), Inches(2.8), left_lines, font="Aptos", size=12, color=WHITE, dark=True)
    add_card(slide, Inches(5.15), Inches(1.8), Inches(6.8), Inches(3.85))
    add_text(slide, Inches(5.45), Inches(2.05), Inches(5.6), Inches(0.45), right_title, font="Georgia", size=18, color=TEXT)
    add_multiline(slide, Inches(5.45), Inches(2.48), Inches(5.9), Inches(2.9), right_lines, font="Aptos", size=12.5, color=MUTED)


def add_stack_architecture(slide):
    set_bg(slide)
    add_header_chrome(slide, "04", "Stack")
    add_card(slide, Inches(0.6), Inches(0.5), Inches(12.1), Inches(2.35))
    add_text(slide, Inches(0.9), Inches(0.82), Inches(4.5), Inches(0.5), "Technology Stack", font="Georgia", size=24, color=TEXT)
    add_text(slide, Inches(0.9), Inches(1.32), Inches(6.3), Inches(0.4), "The project used a modern JavaScript ecosystem focused on frontend performance, developer velocity, and live ecommerce reliability.", font="Aptos", size=12.5, color=MUTED)
    tech_cards = [
        ("Next.js & React", "Server-rendered and client-side UI architecture for storefront and admin experiences"),
        ("JavaScript (ES6+)", "Async logic, reusable components, array methods, event handling, and app flow control"),
        ("Context API", "Shared global state for cart, settings, wishlist, theme, and related frontend behaviors"),
        ("Git & APIs", "Version control, collaboration, RESTful integration, payment flow support, and production fixes"),
    ]
    x_positions = [0.9, 4.2, 7.45, 10.0]
    widths = [2.85, 2.85, 2.25, 2.2]
    for idx, ((title, body), x, w) in enumerate(zip(tech_cards, x_positions, widths), start=1):
        add_icon_badge(slide, Inches(x - 0.03), Inches(1.92), str(idx))
        add_text(slide, Inches(x), Inches(1.95), Inches(w), Inches(0.35), title, font="Georgia", size=15, color=TEXT)
        add_text(slide, Inches(x), Inches(2.3), Inches(w), Inches(0.55), body, font="Aptos", size=10.8, color=MUTED)

    add_card(slide, Inches(0.6), Inches(3.1), Inches(12.1), Inches(3.2))
    add_text(slide, Inches(0.9), Inches(3.42), Inches(5.3), Inches(0.45), "Project Architecture", font="Georgia", size=24, color=TEXT)
    add_text(slide, Inches(0.9), Inches(3.88), Inches(8.2), Inches(0.45), "BogEcom follows a connected product architecture where frontend portals communicate with API and business logic layers serving both customers and administrators.", font="Aptos", size=12.5, color=MUTED)

    center = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(5.3), Inches(4.35), Inches(2.1), Inches(1.5))
    center.fill.solid()
    center.fill.fore_color.rgb = RGBColor(152, 120, 28)
    center.line.color.rgb = RGBColor(152, 120, 28)
    tf = center.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "BogEcom\nCore Platform"
    p.alignment = PP_ALIGN.CENTER
    for run in p.runs:
        run.font.name = "Georgia"
        run.font.size = Pt(16)
        run.font.color.rgb = WHITE

    for diameter in [3.2, 4.2, 5.2]:
        ring = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(5.3) - Inches((diameter - 2.1) / 2), Inches(4.35) - Inches((diameter - 1.5) / 2), Inches(diameter), Inches(diameter * 0.8))
        ring.fill.background()
        ring.line.color.rgb = RGBColor(227, 206, 151)
        ring.line.width = Pt(1.1)

    add_text(slide, Inches(1.05), Inches(4.9), Inches(2.0), Inches(0.35), "Client Portal", font="Georgia", size=16, color=TEXT)
    add_text(slide, Inches(1.05), Inches(5.25), Inches(2.3), Inches(0.5), "Next.js storefront: product pages, cart, checkout, account, membership, orders", font="Aptos", size=11.2, color=MUTED)
    add_text(slide, Inches(9.25), Inches(4.9), Inches(2.0), Inches(0.35), "Admin Portal", font="Georgia", size=16, color=TEXT)
    add_text(slide, Inches(9.25), Inches(5.25), Inches(2.2), Inches(0.5), "Dashboard, orders, settings, combos, analytics, shipping, support workflows", font="Aptos", size=11.2, color=MUTED)


def add_timeline_responsibilities(slide):
    set_bg(slide)
    add_header_chrome(slide, "05", "Timeline")
    add_card(slide, Inches(0.55), Inches(0.5), Inches(12.25), Inches(2.25))
    add_text(slide, Inches(0.9), Inches(0.8), Inches(5.4), Inches(0.45), "Monthly Progress Timeline", font="Georgia", size=24, color=TEXT)
    add_text(slide, Inches(0.9), Inches(1.25), Inches(7.2), Inches(0.35), "Growth followed a deliberate ramp-up from onboarding and exploration to independent feature ownership by the final sprint.", font="Aptos", size=12.2, color=MUTED)
    months = [
        ("January - Foundation", "Onboarding, codebase exploration, environment setup, first UI contributions"),
        ("February - Building", "Page development, Context API integration, cart and product listing features"),
        ("March - Deepening", "Checkout flow, admin panel work, API integration, complex pricing logic"),
        ("April - Ownership", "End-to-end feature delivery, regression testing, bug fixes, deployment support"),
    ]
    xs = [1.45, 4.2, 7.0, 9.8]
    slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(1.4), Inches(2.0), Inches(9.8), Inches(0.02)).line.color.rgb = RGBColor(219, 206, 173)
    for idx, ((title, body), x) in enumerate(zip(months, xs), start=1):
        add_text(slide, Inches(x), Inches(1.48), Inches(2.0), Inches(0.3), title, font="Georgia", size=15, color=TEXT, align=PP_ALIGN.CENTER)
        add_text(slide, Inches(x - 0.2), Inches(1.78), Inches(2.4), Inches(0.45), body, font="Aptos", size=10.6, color=MUTED, align=PP_ALIGN.CENTER)
        badge = add_card(slide, Inches(x + 0.58), Inches(1.85), Inches(0.28), Inches(0.28), fill=RGBColor(248, 237, 202), line=RGBColor(220, 196, 126))
        p = badge.text_frame.paragraphs[0]
        p.text = str(idx)
        p.alignment = PP_ALIGN.CENTER
        for run in p.runs:
            run.font.name = "Georgia"
            run.font.size = Pt(10)
            run.font.color.rgb = TEXT

    add_card(slide, Inches(0.55), Inches(3.05), Inches(12.25), Inches(3.2))
    add_text(slide, Inches(0.9), Inches(3.38), Inches(6.2), Inches(0.45), "Key Responsibilities & Contributions", font="Georgia", size=24, color=TEXT)
    add_text(slide, Inches(1.0), Inches(4.0), Inches(2.6), Inches(0.3), "Core Responsibilities", font="Georgia", size=16, color=TEXT)
    add_multiline(
        slide,
        Inches(1.0),
        Inches(4.35),
        Inches(5.15),
        Inches(1.9),
        [
            "-> Page Development - product listing, detail, cart, checkout pages",
            "-> UI Components - reusable blocks across both portals",
            "-> State Management - cart, settings, wishlist, theme behavior",
            "-> Admin Features - inventory, order management, analytics support",
            "-> API Integration - client endpoints and linked backend behavior",
        ],
        font="Aptos",
        size=12,
        color=MUTED,
        gap=9,
    )
    add_card(slide, Inches(7.1), Inches(4.0), Inches(4.9), Inches(1.95), fill=DARK_PANEL, line=DARK_PANEL)
    add_text(slide, Inches(7.4), Inches(4.18), Inches(4.0), Inches(0.35), "Major Contribution Areas", font="Georgia", size=16, color=WHITE, bold=True)
    add_multiline(
        slide,
        Inches(7.4),
        Inches(4.55),
        Inches(4.15),
        Inches(1.45),
        [
            "Storefront: customer-facing product discovery and purchase journey",
            "UI Components: modular reusable blocks for consistency",
            "State Management: shared patterns powering cart and settings",
            "Admin Portal: internal tools for catalog and order operations",
            "Checkout Flow: end-to-end purchase experience with validation",
        ],
        font="Aptos",
        size=11.3,
        color=WHITE,
        dark=True,
    )


def add_testing_metrics(slide):
    set_bg(slide)
    add_header_chrome(slide, "06", "Execution")
    add_text(slide, Inches(0.85), Inches(0.55), Inches(7.0), Inches(0.45), "Testing, Validation & Deployment", font="Georgia", size=24, color=TEXT)
    add_text(slide, Inches(0.85), Inches(1.0), Inches(7.4), Inches(0.35), "Quality assurance and deployment readiness were treated as integral to every sprint, not as an afterthought.", font="Aptos", size=12.2, color=MUTED)
    cards = [
        ("Manual & Flow-Based Testing", "Executed structured manual test cases across product browsing, cart operations, checkout completion, and admin workflows."),
        ("Regression Awareness", "Maintained awareness of regression risk when modifying shared components and linked business logic."),
        ("Deployment & Stability Support", "Contributed to release readiness by fixing bugs, resolving edge cases, and iterating quickly on reported issues."),
    ]
    x = 0.85
    for idx, (title, body) in enumerate(cards, start=1):
        add_card(slide, Inches(x), Inches(1.45), Inches(3.7), Inches(1.9))
        stripe = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(1.45), Inches(0.06), Inches(1.9))
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = ACCENT_DARK
        stripe.line.width = Pt(0)
        add_icon_badge(slide, Inches(x + 0.12), Inches(1.58), str(idx))
        add_text(slide, Inches(x + 0.15), Inches(1.58), Inches(3.2), Inches(0.45), title, font="Georgia", size=16, color=TEXT)
        add_text(slide, Inches(x + 0.15), Inches(2.05), Inches(3.2), Inches(0.9), body, font="Aptos", size=11.4, color=MUTED)
        x += 3.9

    add_card(slide, Inches(0.85), Inches(3.85), Inches(11.95), Inches(2.45))
    add_text(slide, Inches(1.1), Inches(4.12), Inches(5.5), Inches(0.45), "Work Evidence & Metrics", font="Georgia", size=24, color=TEXT)
    add_text(slide, Inches(1.1), Inches(4.55), Inches(6.0), Inches(0.35), "Quantifiable output across the internship reflects consistent, high-volume contribution to a live production codebase.", font="Aptos", size=12.2, color=MUTED)
    metrics = [
        ("75", "Git Commits", "Consistent, incremental commits across the internship duration"),
        ("351", "Files Modified", "Spanning pages, components, utilities, server logic, and styles"),
        ("75K+", "Lines Added", "Reflecting substantial feature delivery and iterative UI construction"),
        ("4+", "Months Active", "January through April with ongoing latest updates included"),
    ]
    positions = [(1.35, 5.05), (5.05, 5.05), (1.35, 5.9), (5.05, 5.9)]
    for (value, label, note), (x, y) in zip(metrics, positions):
        add_text(slide, Inches(x), Inches(y), Inches(1.8), Inches(0.45), value, font="Georgia", size=24, color=TEXT, align=PP_ALIGN.CENTER)
        add_text(slide, Inches(x), Inches(y + 0.38), Inches(2.2), Inches(0.3), label, font="Georgia", size=14, color=TEXT, align=PP_ALIGN.CENTER)
        add_text(slide, Inches(x - 0.25), Inches(y + 0.7), Inches(2.7), Inches(0.42), note, font="Aptos", size=9.6, color=MUTED, align=PP_ALIGN.CENTER)
    note = add_card(slide, Inches(0.95), Inches(6.9), Inches(11.7), Inches(0.42), fill=BLUE_NOTE, line=RGBColor(153, 188, 236))
    p = note.text_frame.paragraphs[0]
    p.text = "These metrics represent direct contributions to a live ecommerce platform. The report has been updated with latest work through April 21, 2026."
    p.alignment = PP_ALIGN.CENTER
    for run in p.runs:
        run.font.name = "Aptos"
        run.font.size = Pt(10.5)
        run.font.color.rgb = RGBColor(52, 86, 139)


def add_screenshot_gallery(slide, title, left_image, right_image, left_caption, right_caption):
    set_bg(slide)
    add_header_chrome(slide, "07" if title.endswith("I") else "08", "Live Screens")
    add_card(slide, Inches(0.55), Inches(0.48), Inches(12.25), Inches(6.25))
    add_text(slide, Inches(0.9), Inches(0.8), Inches(6.0), Inches(0.45), title, font="Georgia", size=24, color=TEXT)
    add_text(slide, Inches(0.9), Inches(1.22), Inches(8.0), Inches(0.4), "Representative live-site screens showing how implemented work appears in the deployed product experience.", font="Aptos", size=12.2, color=MUTED)
    if left_image.exists():
        slide.shapes.add_picture(str(left_image), Inches(0.95), Inches(1.8), width=Inches(5.45), height=Inches(3.55))
    if right_image.exists():
        slide.shapes.add_picture(str(right_image), Inches(6.55), Inches(1.8), width=Inches(5.25), height=Inches(3.55))
    add_text(slide, Inches(1.05), Inches(5.55), Inches(5.2), Inches(0.35), left_caption, font="Georgia", size=14, color=TEXT, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(6.65), Inches(5.55), Inches(5.0), Inches(0.35), right_caption, font="Georgia", size=14, color=TEXT, align=PP_ALIGN.CENTER)


def add_challenges_conclusion(slide):
    set_bg(slide)
    add_header_chrome(slide, "09", "Close")
    add_text(slide, Inches(0.85), Inches(0.55), Inches(7.0), Inches(0.45), "Challenges, Learning & Conclusion", font="Georgia", size=24, color=TEXT)
    add_card(slide, Inches(0.8), Inches(1.35), Inches(5.2), Inches(4.75))
    add_icon_badge(slide, Inches(1.02), Inches(1.67), "!")
    add_text(slide, Inches(1.1), Inches(1.65), Inches(2.8), Inches(0.4), "Challenges Faced", font="Georgia", size=18, color=TEXT)
    add_multiline(
        slide,
        Inches(1.05),
        Inches(2.08),
        Inches(4.35),
        Inches(3.6),
        [
            "Codebase Complexity: navigating a large existing system with shared components and business logic",
            "Pricing Logic: implementing dynamic calculations with discounts, taxes, variants, and combos",
            "Regression Risk: ensuring changes in shared modules did not break connected flows",
            "Multi-Portal Coordination: balancing client polish, admin behavior, and server-linked correctness",
            "Release Pressure: fixing production-impacting issues while maintaining visual consistency",
        ],
        font="Aptos",
        size=12,
        color=MUTED,
    )
    add_card(slide, Inches(6.3), Inches(1.35), Inches(5.7), Inches(2.3), fill=DARK_PANEL, line=DARK_PANEL)
    add_icon_badge(slide, Inches(6.52), Inches(1.64), "+", dark=True)
    add_text(slide, Inches(6.65), Inches(1.62), Inches(3.4), Inches(0.35), "Learning Outcomes", font="Georgia", size=18, color=WHITE, bold=True)
    add_multiline(
        slide,
        Inches(6.6),
        Inches(2.02),
        Inches(4.8),
        Inches(1.3),
        [
            "Technical: deeper expertise in React, Next.js, Context API, and API integration patterns",
            "Product Thinking: understood how frontend decisions impact user experience and business outcomes",
            "Collaboration: improved confidence in working safely inside an active multi-author repository",
        ],
        font="Aptos",
        size=11.2,
        color=WHITE,
        dark=True,
    )
    add_card(slide, Inches(6.3), Inches(3.95), Inches(5.7), Inches(2.15), fill=RGBColor(249, 239, 210), line=RGBColor(228, 208, 150))
    add_icon_badge(slide, Inches(6.52), Inches(4.22), ">", dark=False)
    add_text(slide, Inches(6.65), Inches(4.2), Inches(4.6), Inches(0.35), "Closing Reflection", font="Georgia", size=18, color=TEXT)
    add_text(
        slide,
        Inches(6.65),
        Inches(4.65),
        Inches(4.7),
        Inches(1.2),
        "This internship gave me real exposure to product-scale frontend engineering. My contribution spanned UI, logic, integration, testing, and delivery support, making the experience a strong bridge between academic learning and professional software work.",
        font="Aptos",
        size=12,
        color=MUTED,
    )


def build() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    home = SCREEN_DIR / "home.png"
    products = SCREEN_DIR / "products.png"
    product_detail = SCREEN_DIR / "product-detail.png"
    membership = SCREEN_DIR / "membership.png"

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_jecrc_title_slide(slide)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_jecrc_contents_slide(slide)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_split_hero(
        slide,
        home,
        "Frontend Developer Internship Report",
        "BogEcom - Buy\nOne Gram",
        "Piyush Songara - B.Tech Computer Science",
        "Internship Guide: Faculty & Industry Mentor - Duration: January - April 2026",
        ["INTERNSHIP REPORT", "FRONTEND DEVELOPMENT", "E-COMMERCE PLATFORM"],
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_overview_slide(slide, membership)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dual_panel(
        slide,
        "Objective & Platform Context",
        "A real product internship focused on delivering production-ready ecommerce experiences across client and admin portals.",
        "Internship Objective",
        [
            "Bridge academic knowledge with real-world software engineering by contributing to a live, production-grade ecommerce product.",
            "Understand frontend architecture at scale.",
            "Deliver features under real product constraints.",
            "Collaborate using Git, APIs, and shared code patterns.",
        ],
        "About BogEcom (Buy One Gram)",
        [
            "BogEcom is a dual-portal ecommerce platform serving both customers and internal administrators.",
            "Client Portal: product browsing, cart, checkout, order tracking.",
            "Admin Portal: inventory, pricing, orders, user management.",
            "Backend: RESTful APIs feeding both portals and linked business workflows.",
        ],
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_stack_architecture(slide)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_timeline_responsibilities(slide)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_testing_metrics(slide)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_screenshot_gallery(slide, "Live Implementation Screens I", home, products, "Home Page on Live Site", "Products Listing on Live Site")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_screenshot_gallery(slide, "Live Implementation Screens II", product_detail, membership, "Product Detail Page on Live Site", "Membership Page on Live Site")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_challenges_conclusion(slide)

    prs.save(str(OUTPUT_PATH))
    print(OUTPUT_PATH)


if __name__ == "__main__":
    build()
