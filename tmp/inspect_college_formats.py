from docx import Document
from pptx import Presentation


def inspect_docx(path: str) -> None:
    doc = Document(path)
    print(f"DOCX: {path}")
    texts = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    for idx, text in enumerate(texts, 1):
        print(f"P{idx}: {text}")
    print(f"TABLES: {len(doc.tables)}")
    for table_index, table in enumerate(doc.tables, 1):
        print(f"TABLE {table_index}: {len(table.rows)}x{len(table.columns)}")
        for row in table.rows:
            print(" | ".join(cell.text.strip() for cell in row.cells))
    print()


def inspect_pptx(path: str) -> None:
    prs = Presentation(path)
    print(f"PPTX: {path}")
    print(f"SLIDES: {len(prs.slides)}")
    for slide_index, slide in enumerate(prs.slides, 1):
        print(f"SLIDE {slide_index}")
        for shape in slide.shapes:
            text = getattr(shape, "text", "").strip()
            if text:
                print(text)
        print()


if __name__ == "__main__":
    inspect_docx(r"d:\BogEcom\Mid Review Report Format.docx")
    inspect_pptx(r"d:\BogEcom\PPT Format.pptx")
    inspect_docx(r"d:\BogEcom\Piyush_Songara_Mid_Review_Report.docx")
    inspect_pptx(r"d:\BogEcom\Piyush_Songara_Mid_Review_Presentation.pptx")
