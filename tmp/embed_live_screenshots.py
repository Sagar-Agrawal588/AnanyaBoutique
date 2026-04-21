from __future__ import annotations

from pathlib import Path

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Inches, Pt
from docx.oxml.ns import qn
from pptx import Presentation
from pptx.util import Inches as PPTInches, Pt as PPTPt
from pptx.enum.text import PP_ALIGN


ROOT = Path(r"d:\BogEcom")
DOCX_PATH = ROOT / "Piyush_Songara_Mid_Review_Report.docx"
PPTX_PATH = ROOT / "Piyush_Songara_Mid_Review_Presentation.pptx"
IMAGE_DIR = ROOT / "output" / "mid-review-live-screenshots"

IMAGES = [
    ("Home Page on Live Site", IMAGE_DIR / "home.png"),
    ("Products Listing on Live Site", IMAGE_DIR / "products.png"),
    ("Product Detail Page on Live Site", IMAGE_DIR / "product-detail.png"),
    ("Membership Page on Live Site", IMAGE_DIR / "membership.png"),
]


def style_run(run, size=12, bold=False):
    run.font.name = "Times New Roman"
    run._element.rPr.rFonts.set(qn("w:eastAsia"), "Times New Roman")
    run.font.size = Pt(size)
    run.font.bold = bold


def add_report_screenshots() -> None:
    doc = Document(str(DOCX_PATH))

    target = None
    for paragraph in doc.paragraphs:
        if paragraph.text.strip() == "5 LEARNING OUTCOME FROM 6 MONTHS INTERNSHIP":
            target = paragraph
            break
    if target is None:
        doc.save(str(DOCX_PATH))
        return

    p = target.insert_paragraph_before()
    p.alignment = WD_ALIGN_PARAGRAPH.LEFT
    run = p.add_run("Live Implementation Screenshots")
    style_run(run, size=14, bold=True)

    p = target.insert_paragraph_before()
    p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
    run = p.add_run(
        "The following screenshots are captured from the live site and show representative pages where my implementation and frontend refinement work is visible in the actual deployed product."
    )
    style_run(run, size=12)

    for caption, image_path in IMAGES:
        if not image_path.exists():
            continue
        cap = target.insert_paragraph_before()
        cap.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = cap.add_run(caption)
        style_run(run, size=12, bold=True)

        pic = target.insert_paragraph_before()
        pic.alignment = WD_ALIGN_PARAGRAPH.CENTER
        pic.add_run().add_picture(str(image_path), width=Inches(5.9))

        desc = target.insert_paragraph_before()
        desc.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
        run = desc.add_run(
            f"This live screenshot documents the deployed appearance of the {caption.lower()}, helping connect the written report with the actual implemented output."
        )
        style_run(run, size=11)

    doc.save(str(DOCX_PATH))


def set_slide_title(shape, text: str):
    shape.text = text
    for paragraph in shape.text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        for run in paragraph.runs:
            run.font.name = "Times New Roman"
            run.font.size = PPTPt(24)
            run.font.bold = True


def add_caption(slide, left, top, width, height, text):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.name = "Times New Roman"
    run.font.size = PPTPt(14)
    run.font.bold = True


def add_screenshot_slides() -> None:
    prs = Presentation(str(PPTX_PATH))

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    set_slide_title(slide.shapes.title, "Live Site Screenshots I")
    slide.shapes.add_picture(str(IMAGES[0][1]), PPTInches(0.6), PPTInches(1.2), width=PPTInches(4.2))
    slide.shapes.add_picture(str(IMAGES[1][1]), PPTInches(5.1), PPTInches(1.2), width=PPTInches(4.2))
    add_caption(slide, PPTInches(0.6), PPTInches(5.8), PPTInches(4.2), PPTInches(0.5), IMAGES[0][0])
    add_caption(slide, PPTInches(5.1), PPTInches(5.8), PPTInches(4.2), PPTInches(0.5), IMAGES[1][0])

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    set_slide_title(slide.shapes.title, "Live Site Screenshots II")
    slide.shapes.add_picture(str(IMAGES[2][1]), PPTInches(0.6), PPTInches(1.2), width=PPTInches(4.2))
    slide.shapes.add_picture(str(IMAGES[3][1]), PPTInches(5.1), PPTInches(1.2), width=PPTInches(4.2))
    add_caption(slide, PPTInches(0.6), PPTInches(5.8), PPTInches(4.2), PPTInches(0.5), IMAGES[2][0])
    add_caption(slide, PPTInches(5.1), PPTInches(5.8), PPTInches(4.2), PPTInches(0.5), IMAGES[3][0])

    prs.save(str(PPTX_PATH))


if __name__ == "__main__":
    add_report_screenshots()
    add_screenshot_slides()
    print(DOCX_PATH)
    print(PPTX_PATH)
