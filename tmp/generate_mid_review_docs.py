from __future__ import annotations

import subprocess
from collections import Counter
from datetime import date
from pathlib import Path

from docx import Document
from docx.enum.section import WD_SECTION
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Inches, Pt
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches as PPTInches
from pptx.util import Pt as PPTPt


ROOT = Path(r"d:\BogEcom")
TEMPLATE_PPT = ROOT / "PPT Format.pptx"
OUTPUT_PPT = ROOT / "Piyush_Songara_Mid_Review_Presentation.pptx"
OUTPUT_DOCX = ROOT / "Piyush_Songara_Mid_Review_Report.docx"
FALLBACK_OUTPUT_DOCX = ROOT / "Piyush_Songara_Mid_Review_Report_Updated.docx"
FALLBACK_OUTPUT_PPT = ROOT / "Piyush_Songara_Mid_Review_Presentation_Updated.pptx"
LOGO_PATH = ROOT / "frontend" / "client" / "public" / "logo.png"
REPORT_PATHS = ["frontend/client", "frontend/admin", "server"]


def run_git(*args: str) -> str:
    result = subprocess.run(
        ["git", *args],
        cwd=ROOT,
        check=True,
        capture_output=True,
        text=True,
        encoding="utf-8",
    )
    return result.stdout


def get_repo_stats() -> dict[str, object]:
    email = "piyushsongara69@gmail.com"
    commit_count = int(
        run_git("rev-list", "--count", f"--author={email}", "HEAD", "--", *REPORT_PATHS).strip()
    )

    file_lines = run_git(
        "log",
        f"--author={email}",
        "--name-only",
        "--pretty=format:",
        "--",
        *REPORT_PATHS,
    ).splitlines()
    files = [line.strip() for line in file_lines if line.strip()]
    unique_files = sorted(set(files))
    file_counter = Counter(files)

    dates = [
        line.strip()
        for line in run_git(
            "log",
            f"--author={email}",
            "--format=%ad",
            "--date=short",
            "--",
            *REPORT_PATHS,
        ).splitlines()
        if line.strip()
    ]

    recent_commits = [
        line.strip()
        for line in run_git(
            "log",
            "--oneline",
            f"--author={email}",
            "-n",
            "10",
            "--",
            *REPORT_PATHS,
        ).splitlines()
        if line.strip()
    ]

    authored_commit_entries = []
    for line in run_git(
        "log",
        f"--author={email}",
        "--date=short",
        '--pretty=format:%ad|%h|%s',
        "--",
        *REPORT_PATHS,
    ).splitlines():
        line = line.strip().strip('"')
        if not line:
            continue
        parts = line.split("|", 2)
        if len(parts) == 3:
            authored_commit_entries.append(
                {"date": parts[0], "hash": parts[1], "subject": parts[2]}
            )

    overall_shortlog = run_git("shortlog", "-sne", "--all").splitlines()
    overall_commits = 0
    for line in overall_shortlog:
        if email in line:
            overall_commits += int(line.strip().split()[0])

    top_files = file_counter.most_common(8)

    category_counter: Counter[str] = Counter()
    for file_path in unique_files:
        normalized = file_path.replace("\\", "/")
        if normalized.startswith("server/"):
            category_counter["Server and API"] += 1
        elif "/src/components/" in normalized:
            category_counter["UI components"] += 1
        elif "/src/app/" in normalized:
            category_counter["Pages and flows"] += 1
        elif "/src/context/" in normalized:
            category_counter["State management"] += 1
        elif "/src/hooks/" in normalized:
            category_counter["Hooks and integration"] += 1
        elif "/src/utils/" in normalized:
            category_counter["Utilities and API"] += 1
        else:
            category_counter["Config and assets"] += 1

    monthly_commit_counter: Counter[str] = Counter()
    for entry in authored_commit_entries:
        monthly_commit_counter[entry["date"][:7]] += 1

    numstat_lines = run_git(
        "log",
        f"--author={email}",
        "--numstat",
        "--pretty=tformat:",
        "--",
        *REPORT_PATHS,
    ).splitlines()
    churn_counter: Counter[str] = Counter()
    additions = 0
    deletions = 0
    for line in numstat_lines:
        parts = line.split("\t")
        if len(parts) != 3:
            continue
        added, deleted, file_path = parts
        if added != "-":
            additions += int(added)
        if deleted != "-":
            deletions += int(deleted)
        if added != "-" and deleted != "-":
            churn_counter[file_path] += int(added) + int(deleted)

    monthly_breakdown = [
        (
            "2026-01",
            "January 2026",
            monthly_commit_counter.get("2026-01", 0),
            "Started contributing to foundational storefront work including authentication, payment hooks, layout, landing page modules, and theme-driven UI behavior.",
        ),
        (
            "2026-02",
            "February 2026",
            monthly_commit_counter.get("2026-02", 0),
            "Worked heavily on checkout, address, orders, membership, support, notifications, and broader ecommerce flow fixes across both client and admin surfaces.",
        ),
        (
            "2026-03",
            "March 2026",
            monthly_commit_counter.get("2026-03", 0),
            "Focused on combo flows, product cards, zoom/image experience, wishlist, settings, shipping, and UI consistency improvements across multiple screens.",
        ),
        (
            "2026-04",
            "April 2026",
            monthly_commit_counter.get("2026-04", 0),
            "Closed the period with order-id flow fixes, responsive product-detail improvements, share flow updates, payment handling updates, admin authentication improvements, and final bug-fix work across connected modules.",
        ),
    ]

    return {
        "student_name": "Piyush Songara",
        "program": "Bachelor of Technology in Computer Science and Engineering",
        "semester": "Internship Semester Jan-June 2026",
        "company_name": "BogEcom / Buy One Gram E-commerce Platform",
        "role": "Frontend Developer Intern",
        "department": "Computer Science and Engineering",
        "institute": "JECRC University, Jaipur",
        "faculty_guide": "[Update Faculty Guide Name]",
        "industry_guide": "[Update Industry Guide Name]",
        "registration_number": "[Update Registration Number]",
        "submission_date": "April 2026",
        "start_date": min(dates) if dates else "2026-01-30",
        "end_date": max(dates) if dates else "2026-04-05",
        "frontend_commit_count": commit_count,
        "overall_commit_count": overall_commits,
        "unique_files": len(unique_files),
        "unique_files_list": unique_files,
        "top_files": top_files,
        "all_file_touch_counts": file_counter.most_common(),
        "top_churn_files": churn_counter.most_common(20),
        "recent_commits": recent_commits,
        "authored_commit_entries": authored_commit_entries,
        "category_counter": category_counter,
        "monthly_breakdown": monthly_breakdown,
        "total_additions": additions,
        "total_deletions": deletions,
        "work_summary": [
            "Implemented and refined core storefront pages including product details, cart, checkout, orders, membership, address, and account flows.",
            "Built and improved reusable UI modules such as Header, Footer, ProductItem, ProductSlider, CartDrawer, Search, FlavorSwitcherBar, and support views.",
            "Fixed pricing, GST, coupon, combo, order-display, and checkout logic by aligning frontend calculations with backend behavior.",
            "Integrated cart state, wishlist state, notifications, share handling, support flows, and payment-related updates across the client app.",
            "Contributed to admin-side pages such as orders, settings, shipping, customer care, purchase orders, analytics, and membership management.",
            "Handled API and payment-related changes in recent work, including PhonePe flow refinement, authorization logic updates, and linked order context fixes.",
        ],
        "outcomes": [
            "Improved responsiveness and consistency of product listing and product detail experiences.",
            "Reduced checkout and cart errors by fixing order calculation and pricing edge cases.",
            "Improved user-facing flows for membership, support, notifications, and order tracking.",
            "Maintained both client and admin interfaces while coordinating with backend-driven API changes.",
        ],
        "tools": [
            "Next.js",
            "React.js",
            "JavaScript",
            "Context API",
            "CSS / Tailwind-style utility classes",
            "Git and GitHub",
            "REST API integration",
            "VS Code",
        ],
        "methodology": [
            "Understand the bug or feature requirement from existing flow and affected screens.",
            "Trace connected components, pages, context state, and API utilities.",
            "Implement UI updates and behavior changes in reusable frontend modules.",
            "Verify cart, checkout, orders, and other impacted user journeys locally.",
            "Refine responsiveness, edge-case handling, and cross-page consistency before commit.",
        ],
        "monthly_progress": [
            (
                "January 2026",
                "Onboarded into the codebase, worked on authentication, payment hooks, theme handling, landing page components, and initial storefront UI stabilization.",
            ),
            (
                "February 2026",
                "Handled membership, address, cart, checkout, search, customer care, notification, and order-related frontend fixes across both client and admin.",
            ),
            (
                "March 2026",
                "Focused on combo deals, product cards, zoom and image handling, wishlist, newsletter, settings, shipping, and broader ecommerce flow improvements.",
            ),
            (
                "April 2026",
                "Delivered order-id flow fixes, responsive product-detail refinements, share flow updates, and final polishing of product display components.",
            ),
        ],
        "learning_points": [
            "Learned to maintain a production-style Next.js codebase with shared state and multi-page ecommerce flows.",
            "Improved debugging skills for pricing, cart, order, and checkout-related edge cases.",
            "Gained experience in balancing UI design, frontend logic, and backend API integration.",
            "Developed discipline in working with an existing team codebase and iterative Git-based delivery.",
        ],
        "detailed_role_areas": [
            (
                "Storefront Pages",
                "A major portion of my internship involved improving pages that directly influence the customer purchase journey. These included product detail screens, category and product listing pages, cart, checkout, my orders, order detail, membership, address management, and account-related pages. The work required attention to layout consistency, state synchronization, dynamic data rendering, and bug handling under different user scenarios.",
            ),
            (
                "Reusable UI Components",
                "I repeatedly worked on shared components such as Header, Footer, ProductItem, ProductSlider, PopularProducts, CartDrawer, Search, FlavorSwitcherBar, ComboCard, ProductZoom, and related product detail widgets. These components were central to maintaining consistency across the storefront and were often the first point for fixing responsiveness, rendering issues, or interaction problems.",
            ),
            (
                "State Management and Flow Logic",
                "A significant part of my contribution extended beyond visible UI changes into state management through Context-based modules such as CartContext, WishlistContext, SettingsContext, ThemeContext, and related helpers. This work was especially important in cart quantity management, combo handling, order totals, coupon behavior, and synchronization between pages that consume shared data.",
            ),
            (
                "Admin Support Screens",
                "Although my strongest contribution area was the client frontend, I also worked on admin-facing pages such as orders, settings, shipping, purchase orders, customer care, membership, products, and layout/sidebar structures. This gave me broader understanding of how operational workflows connect with the storefront experience.",
            ),
            (
                "Server and Integration Support",
                "In the more recent phase of work, my contribution also included connected API and payment-flow changes. This involved handling logic around payment processing, order context, and authorization-related updates that influenced how frontend and backend behavior stayed aligned.",
            ),
        ],
        "case_studies": [
            (
                "Checkout and Pricing Stability",
                "Checkout-related fixes formed one of the most important areas of my internship because errors in totals, GST, coupons, combo pricing, or order presentation can directly impact user trust and payment correctness. My work included aligning frontend calculations with backend rules, reducing pricing mismatches, and improving overall reliability of cart and checkout transitions.",
            ),
            (
                "Product Card and Product Detail Improvements",
                "Product presentation is central to ecommerce conversion, so I spent substantial time on ProductItem, ProductSlider, PopularProducts, and product detail screens. These tasks involved responsiveness, variant-aware display behavior, image handling, badges, and card-level interaction patterns that affect how users discover and compare products.",
            ),
            (
                "Account, Orders, and Support Experience",
                "Beyond discovery and checkout, I contributed to pages that support post-purchase confidence, including my orders, order detail, address pages, contact/support flows, and account interfaces. This work improved continuity of the user experience and reinforced the idea that frontend quality extends far beyond the landing page.",
            ),
            (
                "Membership and Feature Expansion",
                "Membership-related pages and exclusive experience elements were another recurring area. I worked on membership screens and related UI behavior, which exposed me to feature-specific journeys where visual experience, state management, and business logic all need to stay aligned.",
            ),
        ],
        "challenges": [
            "Understanding an existing codebase with interconnected pages, shared contexts, and utility layers required careful reading before making safe modifications.",
            "Pricing-sensitive flows such as checkout, combo selection, GST, discounts, and order totals demanded high correctness because small frontend errors could affect business logic and user trust.",
            "Maintaining responsiveness across reusable components was challenging because visual changes in shared elements could affect multiple screens simultaneously.",
            "Working across both client and admin surfaces required balancing UI polish, flow correctness, and compatibility with backend-driven data contracts.",
        ],
        "testing_observations": [
            "Manual validation remained essential for verifying real user journeys such as add-to-cart, checkout, order viewing, support submission, and membership interaction.",
            "Shared components required regression checking because a change in one module could affect several pages using the same component.",
            "Pricing and order workflows demanded scenario-based testing with different data combinations such as combos, discount behavior, stock display, and variant selection.",
            "Responsive behavior needed repeated checks because many of the improvements directly affected mobile-friendly card layouts, headers, drawers, and detail pages.",
        ],
    }


def set_doc_defaults(doc: Document) -> None:
    style = doc.styles["Normal"]
    style.font.name = "Times New Roman"
    style._element.rPr.rFonts.set(qn("w:eastAsia"), "Times New Roman")
    style.font.size = Pt(12)

    for section in doc.sections:
        section.top_margin = Inches(1)
        section.bottom_margin = Inches(1)
        section.right_margin = Inches(1)
        section.left_margin = Inches(1.25)


def add_page_number(paragraph):
    run = paragraph.add_run()
    fld_char1 = OxmlElement("w:fldChar")
    fld_char1.set(qn("w:fldCharType"), "begin")
    instr_text = OxmlElement("w:instrText")
    instr_text.set(qn("xml:space"), "preserve")
    instr_text.text = "PAGE"
    fld_char2 = OxmlElement("w:fldChar")
    fld_char2.set(qn("w:fldCharType"), "end")
    run._r.append(fld_char1)
    run._r.append(instr_text)
    run._r.append(fld_char2)


def add_heading(doc: Document, text: str, size: int = 16) -> None:
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run(text.upper())
    run.bold = True
    run.font.name = "Times New Roman"
    run.font.size = Pt(size)


def add_subheading(doc: Document, text: str) -> None:
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.LEFT
    run = p.add_run(text)
    run.bold = True
    run.font.name = "Times New Roman"
    run.font.size = Pt(14)


def add_body(doc: Document, text: str, bold: bool = False, align=WD_ALIGN_PARAGRAPH.JUSTIFY) -> None:
    p = doc.add_paragraph()
    p.alignment = align
    p.paragraph_format.line_spacing = 1.5
    run = p.add_run(text)
    run.bold = bold
    run.font.name = "Times New Roman"
    run.font.size = Pt(12)


def add_bullet(doc: Document, text: str) -> None:
    p = doc.add_paragraph(style="List Bullet")
    p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
    p.paragraph_format.line_spacing = 1.5
    run = p.add_run(text)
    run.font.name = "Times New Roman"
    run.font.size = Pt(12)


def add_table_of_contents_placeholder(doc: Document) -> None:
    p = doc.add_paragraph()
    run = p.add_run()
    fld_char1 = OxmlElement("w:fldChar")
    fld_char1.set(qn("w:fldCharType"), "begin")
    instr_text = OxmlElement("w:instrText")
    instr_text.set(qn("xml:space"), "preserve")
    instr_text.text = 'TOC \\o "1-3" \\h \\z \\u'
    fld_char2 = OxmlElement("w:fldChar")
    fld_char2.set(qn("w:fldCharType"), "separate")
    fld_char3 = OxmlElement("w:fldChar")
    fld_char3.set(qn("w:fldCharType"), "end")
    run._r.append(fld_char1)
    run._r.append(instr_text)
    run._r.append(fld_char2)
    run._r.append(fld_char3)


def add_index_table(doc: Document) -> None:
    table = doc.add_table(rows=1, cols=3)
    table.style = "Table Grid"
    headers = ["Sr no.", "Description", "Page no."]
    for idx, value in enumerate(headers):
        cell = table.rows[0].cells[idx]
        cell.text = value
        for paragraph in cell.paragraphs:
            paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
            for run in paragraph.runs:
                run.bold = True
                run.font.name = "Times New Roman"
                run.font.size = Pt(12)

    rows = [
        ("i", "DECLARATION", "i"),
        ("ii", "ACKNOWLEDGEMENT", "ii"),
        ("iii", "INDEX", "iii"),
        ("iv", "LIST OF ILLUSTRATIONS", "iv"),
        ("v", "ABSTRACT", "1"),
        ("1", "INTRODUCTION", "2"),
        ("2", "INTRODUCTION ABOUT COMPANY", "4"),
        ("3", "ABOUT PROJECT", "6"),
        ("3.1", "TECHNOLOGIES WORKED ON", "7"),
        ("3.2", "MONTHLY PROGRESSION OF PROJECT", "9"),
        ("3.3", "TESTING (IF DONE)", "10"),
        ("3.4", "DEPLOYMENT (IF DONE)", "13"),
        ("4", "ANY OTHER POINTS", "14"),
        ("5", "LEARNING OUTCOME FROM 6 MONTHS INTERNSHIP", "18"),
        ("6", "CONCLUSION", "20"),
        ("", "REFERENCES", "21"),
    ]
    for sr_no, desc, page_no in rows:
        row = table.add_row().cells
        row[0].text = sr_no
        row[1].text = desc
        row[2].text = page_no


def add_illustrations_table(doc: Document) -> None:
    table = doc.add_table(rows=1, cols=3)
    table.style = "Table Grid"
    headers = ["Illustration Type", "Illustration Description", "Page Number"]
    for idx, value in enumerate(headers):
        cell = table.rows[0].cells[idx]
        cell.text = value
        for paragraph in cell.paragraphs:
            paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
            for run in paragraph.runs:
                run.bold = True
                run.font.name = "Times New Roman"
                run.font.size = Pt(12)

    rows = [
        ("Table", "", ""),
        ("Table 1", "Index Table", "iii"),
        ("Table 2", "List of Illustrations", "iv"),
        ("Table 3", "Monthly Progression of Project", "9"),
        ("Table 4", "Contribution Summary by Work Area", "15"),
        ("Figure", "", ""),
        ("Figure 1", "Project Architecture Overview", "5"),
        ("Figure 2", "Frontend Work Distribution", "14"),
        ("Figure 3", "Testing and Validation Flow", "12"),
        ("Figure 4", "Learning Outcome Framework", "19"),
    ]
    for kind, desc, page in rows:
        row = table.add_row().cells
        row[0].text = kind
        row[1].text = desc
        row[2].text = page


def build_report(data: dict[str, object]) -> None:
    doc = Document()
    set_doc_defaults(doc)

    add_heading(doc, "MID REVIEW INTERNSHIP REPORT")
    add_body(doc, f"({data['semester']})", align=WD_ALIGN_PARAGRAPH.CENTER)
    add_body(
        doc,
        "A report submitted in partial fulfillment of the requirements for the award of degree of",
        align=WD_ALIGN_PARAGRAPH.CENTER,
    )
    add_heading(doc, "BACHELOR OF TECHNOLOGY", 14)
    add_body(doc, "in", align=WD_ALIGN_PARAGRAPH.CENTER)
    add_body(doc, "COMPUTER SCIENCE AND ENGINEERING", bold=True, align=WD_ALIGN_PARAGRAPH.CENTER)
    add_body(doc, "By", align=WD_ALIGN_PARAGRAPH.CENTER)
    add_body(doc, str(data["student_name"]), bold=True, align=WD_ALIGN_PARAGRAPH.CENTER)
    add_body(doc, f"Reg. No. {data['registration_number']}", align=WD_ALIGN_PARAGRAPH.CENTER)
    if LOGO_PATH.exists():
        doc.add_picture(str(LOGO_PATH), width=Inches(1.8))
        doc.paragraphs[-1].alignment = WD_ALIGN_PARAGRAPH.CENTER
    add_body(doc, "Under the Guidance of", bold=True, align=WD_ALIGN_PARAGRAPH.CENTER)
    add_body(
        doc,
        f"Faculty Internship Guide: {data['faculty_guide']}            Industry Guide: {data['industry_guide']}",
        align=WD_ALIGN_PARAGRAPH.CENTER,
    )
    add_body(doc, str(data["institute"]), bold=True, align=WD_ALIGN_PARAGRAPH.CENTER)
    add_body(doc, str(data["submission_date"]), align=WD_ALIGN_PARAGRAPH.CENTER)

    doc.add_page_break()
    add_heading(doc, "DECLARATION")
    add_body(
        doc,
        f"I hereby declare that the internship work carried out on the project \"{data['company_name']}\" is an authentic record of my work completed as part of the six-month internship requirement for the award of B.Tech in Computer Science and Engineering. The work was performed under the guidance of {data['industry_guide']} and {data['faculty_guide']}.",
    )
    add_body(doc, "(Sign of Student)", align=WD_ALIGN_PARAGRAPH.LEFT)
    add_body(doc, f"Name of Student: {data['student_name']}", align=WD_ALIGN_PARAGRAPH.LEFT)
    add_body(doc, f"RN No.: {data['registration_number']}", align=WD_ALIGN_PARAGRAPH.LEFT)
    add_body(doc, "Place: JECRC University, Jaipur", align=WD_ALIGN_PARAGRAPH.LEFT)
    add_body(doc, f"Date: {date.today().strftime('%d %B %Y')}", align=WD_ALIGN_PARAGRAPH.LEFT)

    doc.add_page_break()
    add_heading(doc, "ACKNOWLEDGEMENT")
    add_body(
        doc,
        f"I would like to express my sincere gratitude to {data['faculty_guide']} for academic guidance and to {data['industry_guide']} for technical mentorship during my internship work on {data['company_name']}. Their continuous support helped me understand real project workflows, improve my frontend engineering skills, and complete the assigned work responsibly.",
    )
    add_body(
        doc,
        "I am also thankful to the department, institute, and project team members who provided a collaborative environment for learning, implementation, testing, and delivery. This internship gave me practical exposure to real ecommerce product development, issue fixing, and team-based software iteration.",
    )
    add_body(doc, "Thanking You.", align=WD_ALIGN_PARAGRAPH.LEFT)
    add_body(doc, str(data["student_name"]), bold=True, align=WD_ALIGN_PARAGRAPH.LEFT)

    doc.add_page_break()
    add_heading(doc, "INDEX")
    add_index_table(doc)

    doc.add_page_break()
    add_heading(doc, "LIST OF ILLUSTRATIONS")
    add_body(
        doc,
        "The list of illustrations gives systematic information about the tables and figures used in this report. It provides a structured overview of where important visual summaries and chapter support material appear in the document.",
    )
    add_illustrations_table(doc)

    doc.add_page_break()
    add_heading(doc, "ABSTRACT")
    add_body(
        doc,
        f"This report presents the work completed during my internship on {data['company_name']}, where I contributed primarily as a {data['role']}. The project is an ecommerce platform with separate client and admin interfaces supported by connected server logic. My contribution focused mainly on frontend development, UI refinement, state management, order and checkout flow improvements, and integration-oriented updates across user-facing and operational journeys.",
    )
    add_body(
        doc,
        f"Using repository history as evidence, I contributed to approximately {data['frontend_commit_count']} authored commits across the covered project modules, spanning about {data['unique_files']} unique files across frontend, admin, and linked server work between {data['start_date']} and {data['end_date']}. The cumulative code churn visible in my authored history covers roughly {data['total_additions']} added lines and {data['total_deletions']} deleted lines, which indicates sustained iteration rather than isolated edits. This internship strengthened my understanding of production-grade React and Next.js development, reusable component design, debugging of complex pricing workflows, and collaborative software delivery.",
    )
    add_body(
        doc,
        "The report not only summarizes my functional responsibilities, but also documents the progression of work through repository-backed evidence, month-wise development patterns, representative files, and detailed case-based explanations of the types of problems I solved. The objective is to present a realistic view of internship contribution in an active product environment rather than a generic overview of software development concepts.",
    )

    add_heading(doc, "1 INTRODUCTION")
    add_body(
        doc,
        "The internship was focused on practical software development in a real project codebase. Instead of working on isolated exercises, I contributed to an active ecommerce application that required continuous feature enhancement, issue resolution, and UI consistency improvements across multiple user journeys.",
    )
    add_body(
        doc,
        "My role centered on frontend implementation. This involved understanding the interaction between pages, reusable components, API utilities, and global application state, then making code changes that improved usability, correctness, and responsiveness.",
    )
    add_body(
        doc,
        "From the beginning of the internship, the work demanded more than visual page changes. Even a seemingly small frontend issue often involved tracing state across shared contexts, checking utility functions, validating API data shapes, and testing how multiple pages behaved after the change. This made the internship highly valuable from an engineering perspective because it exposed me to the practical complexity of maintaining a live, evolving product.",
    )
    add_body(
        doc,
        "The project also helped me understand the relationship between frontend engineering and product quality. In an ecommerce system, component-level mistakes can affect discoverability, conversion, order confidence, and user trust. Therefore, my work often balanced aesthetics, responsiveness, state correctness, and business-rule alignment in the same change set.",
    )
    add_body(
        doc,
        "This report documents the nature of that work. It covers the project environment, technologies used, month-wise progress, major contribution areas, representative problem-solving examples, lessons learned, and repository evidence that supports the summary of effort.",
    )

    add_heading(doc, "2 INTRODUCTION ABOUT COMPANY")
    add_body(
        doc,
        f"{data['company_name']} is an ecommerce-oriented software project with separate user-facing and admin-facing interfaces. The platform supports product browsing, combo deals, membership features, cart and checkout flows, orders, support, content sections, and various admin management panels.",
    )
    add_body(
        doc,
        "The project structure reflects a realistic product environment where frontend logic must remain aligned with backend APIs, business rules, and edge-case handling. This created a meaningful environment for learning how production issues are diagnosed and resolved.",
    )
    add_body(
        doc,
        "Figure 1: Project Architecture Overview. The project can be understood as a linked system of customer storefront, admin dashboard, shared frontend logic, and backend-connected APIs that together support product browsing, order placement, and operational workflows.",
    )
    add_body(
        doc,
        "The client application is responsible for the customer experience. It includes the home page, category browsing, product details, cart, checkout, order history, account settings, membership-specific sections, and support-related pages. Each of these surfaces contributes to the overall usability and conversion quality of the platform.",
    )
    add_body(
        doc,
        "The admin application supports internal workflows such as order management, shipping configuration, purchase orders, customer care, product management, and settings. My exposure to these areas helped me understand how administrative tools shape the data and behavior ultimately seen by the customer-facing frontend.",
    )
    add_body(
        doc,
        "Because both interfaces are part of the same broader system, the internship provided a strong opportunity to learn cross-functional thinking. Many frontend updates made more sense only when viewed in the context of the backend contract, operational requirement, or user journey they were meant to support.",
    )

    add_heading(doc, "3 ABOUT PROJECT")
    add_body(
        doc,
        "The project consists of two major frontend applications: a client storefront for end users and an admin dashboard for management and operational tasks. My work was concentrated around customer-facing flows such as product cards, product details, cart, checkout, orders, membership, account screens, notifications, search, support, and complementary admin views connected to the same business workflows.",
    )
    add_body(
        doc,
        "Technically, the system is organized around Next.js pages, reusable React components, shared context state, utility modules, service integrations, and media-heavy product presentation. This architecture encourages reuse, but it also means changes to common building blocks can have broad effects across several routes.",
    )
    add_body(
        doc,
        "My role in this project was not limited to implementing new screens. A large share of the work involved understanding existing behavior, identifying the correct place for change, minimizing regressions, and preserving consistency between pages that reuse the same data or component layers.",
    )
    add_subheading(doc, "3.1 TECHNOLOGIES WORKED ON")
    add_body(
        doc,
        "The following tools and technologies were used directly or indirectly during the internship. These technologies shaped the implementation approach, debugging workflow, and deployment readiness of frontend features and fixes.",
    )
    for tool in data["tools"]:
        add_bullet(doc, str(tool))
    add_body(
        doc,
        "Among these, Next.js and React formed the primary implementation environment. Context-based state management was especially important because many issues involved cart behavior, settings, themes, wishlist status, or shared user-facing flows. Git-based version control made it possible to track progress, understand change history, and work within a collaborative repository.",
    )

    add_subheading(doc, "3.2 MONTHLY PROGRESSION OF PROJECT")
    table = doc.add_table(rows=1, cols=2)
    table.style = "Table Grid"
    table.rows[0].cells[0].text = "Month"
    table.rows[0].cells[1].text = "Progress Summary"
    for month, summary in data["monthly_progress"]:
        row = table.add_row().cells
        row[0].text = month
        row[1].text = summary
    add_body(
        doc,
        "The table above gives a compact summary, but the actual progression of work was richer and increasingly interconnected over time. As I became more familiar with the codebase, the nature of my work shifted from simpler orientation-level tasks to deeper fixes spanning multiple modules and business-sensitive flows.",
    )
    for key, month, count, summary in data["monthly_breakdown"]:
        add_subheading(doc, f"{month} Detailed Progress")
        add_body(
            doc,
            f"Repository history shows {count} authored commits during {month} across the covered project modules. This indicates active participation across the internship period and shows how my responsibilities evolved with project familiarity.",
        )
        add_body(doc, summary)
        if key == "2026-01":
            add_body(
                doc,
                "In the initial phase, I focused on understanding project structure, page organization, and the way shared contexts affected different parts of the storefront. Early contributions around authentication, payment hooks, landing page behavior, and theme-based features gave me a practical introduction to the codebase.",
            )
            add_body(
                doc,
                "This month established the foundation for later work. By touching layout, home page, and common components early, I built enough confidence to start contributing to more sensitive user flows in the following months.",
            )
        elif key == "2026-02":
            add_body(
                doc,
                "February was the heaviest month in terms of breadth. The commit history from this period reflects work across checkout, address forms, membership, customer care, notifications, order tracking, settings, search, and supporting admin screens. This phase demanded strong understanding of both page composition and underlying utility/state behavior.",
            )
            add_body(
                doc,
                "Several of these tasks were operationally important because they affected the correctness of order and pricing flows. The work during this month significantly improved my ability to debug state propagation, data dependencies, and page-to-page continuity.",
            )
        elif key == "2026-03":
            add_body(
                doc,
                "During March, my contributions became more centered on refinement and feature maturity. Combo-related UX, product card behavior, zoom and image presentation, wishlist integration, settings adjustments, shipping support, and UI consistency improvements became recurring themes.",
            )
            add_body(
                doc,
                "This stage highlighted the importance of reusable component quality. Changes in shared product cards, sliders, drawers, and detail widgets had visible impact across multiple screens, which made regression awareness and component-level thinking especially important.",
            )
        elif key == "2026-04":
            add_body(
                doc,
                "The final stretch of the reviewed period included order-id flow fixes, responsive improvements for product presentation, image and share-related updates, and final polishing of product display elements. These are representative of the kind of finishing work often required before a smoother release experience.",
            )
            add_body(
                doc,
                "By this point, my contribution style was less about orientation and more about targeted refinement of production-facing issues. This reflects how the internship steadily transitioned me into more independent and context-aware frontend work.",
            )

    add_subheading(doc, "3.3 TESTING")
    add_body(
        doc,
        "A large part of the testing during this internship was flow-based validation. After code changes, I verified whether affected journeys such as cart updates, checkout totals, product detail rendering, order pages, wishlist behavior, support pages, and admin panels behaved correctly.",
    )
    add_body(
        doc,
        "This testing approach required checking both normal and edge conditions, such as combo pricing, order identifier rendering, notification behavior, stock display, mobile responsiveness, and consistency of API-driven content.",
    )
    add_body(
        doc,
        "Because the codebase includes several reusable layers, manual validation was rarely limited to the exact file that changed. For example, a modification in CartContext or a shared product component could influence cart, checkout, product detail, and order pages at the same time. This required deliberate regression awareness even when a change request initially looked local.",
    )
    add_body(
        doc,
        "Another important aspect of testing was business-flow correctness. In ecommerce applications, frontend behavior must align with backend expectations around totals, discounts, stock, and identifiers. My work repeatedly involved verifying that rendered values and user-facing behavior matched the intended system logic.",
    )
    add_body(
        doc,
        "Figure 3: Testing and Validation Flow. In practice, testing followed a cycle of requirement understanding, source tracing, implementation, manual flow validation, regression checking, and readiness review before the work was treated as stable.",
    )
    for point in data["testing_observations"]:
        add_bullet(doc, str(point))

    add_subheading(doc, "3.4 DEPLOYMENT / DELIVERY READINESS")
    add_body(
        doc,
        "Although deployment responsibilities were shared at team level, my contribution supported delivery readiness by stabilizing UI behavior, fixing production-impacting issues, and aligning frontend behavior with backend data contracts. This reduced friction during release preparation and regression handling.",
    )
    add_body(
        doc,
        "This contribution is important because frontend delivery readiness is not only about building screens. It also involves ensuring that pages remain stable under real data, common edge cases are handled, and the experience remains consistent across connected flows. Many of my fixes directly supported that readiness by reducing mismatch, breakage, or poor user experience before changes moved further in the pipeline.",
    )

    add_heading(doc, "4 ANY OTHER POINTS")
    add_subheading(doc, "Evidence of Work from Repository History")
    add_bullet(doc, f"Authored commits across covered project modules: {data['frontend_commit_count']}")
    add_bullet(doc, f"Overall commits in repository across my author identities: {data['overall_commit_count']}")
    add_bullet(doc, f"Unique files touched across frontend, admin, and server: {data['unique_files']}")
    add_bullet(doc, f"Approximate authored line additions: {data['total_additions']}")
    add_bullet(doc, f"Approximate authored line deletions: {data['total_deletions']}")
    for label, count in data["category_counter"].most_common():
        add_bullet(doc, f"{label}: {count} unique files")
    add_body(
        doc,
        "These metrics do not claim sole ownership of the project, but they do provide strong evidence of sustained and broad contribution. The spread across pages, components, contexts, utilities, admin screens, and connected server logic supports the conclusion that my work was substantial, recurring, and strongly implementation-oriented with frontend as the major focus.",
    )
    add_body(
        doc,
        "Figure 2: Frontend Work Distribution. The contribution pattern was strongest in pages and flows, followed by reusable UI components, utilities, settings/state handling, and admin support areas. This distribution reflects a frontend role centered on user journey quality rather than only isolated styling updates.",
    )

    add_subheading(doc, "Key Responsibilities Executed")
    for point in data["work_summary"]:
        add_bullet(doc, str(point))
    add_body(
        doc,
        "Collectively, these responsibilities show that my internship work operated at three levels simultaneously: customer-facing UI quality, logic correctness for business-sensitive flows, and maintainability of reusable frontend building blocks.",
    )

    add_subheading(doc, "Detailed Contribution Areas")
    for title, paragraph in data["detailed_role_areas"]:
        add_body(doc, f"{title}: {paragraph}")

    add_subheading(doc, "Project Impact of My Work")
    add_body(
        doc,
        "My work had visible impact on both usability and correctness. On the client side, many changes improved how products were displayed, how users moved through cart and checkout, and how order-related details were shown after purchase. On the admin side, the work supported smoother operational interfaces for orders, shipping, purchase orders, combos, and settings.",
    )
    add_body(
        doc,
        "The most valuable aspect of this contribution was that several changes affected high-trust areas of the platform. In ecommerce applications, issues in totals, product presentation, order identifiers, or account-related information can directly reduce user confidence. Therefore, many of the fixes carried more significance than simple UI polishing.",
    )

    add_subheading(doc, "Coordination Between Client and Admin Work")
    add_body(
        doc,
        "One important learning from the project was that storefront and admin work cannot be treated as separate worlds. Many client-side changes made more sense only after understanding what the admin panel controlled or how operational data was generated. This coordination improved my awareness of the full product lifecycle from internal configuration to customer-facing experience.",
    )
    add_body(
        doc,
        "For example, order flows, shipping information, product management, combo handling, and membership-related experiences are shaped by both backend and admin configurations. Contributing to both sides of these flows helped me build a more practical understanding of feature ownership and data dependency.",
    )

    add_subheading(doc, "Representative Files Frequently Touched")
    for file_path, count in data["top_files"]:
        add_bullet(doc, f"{file_path} ({count} touches in authored commits)")
    add_body(
        doc,
        "The distribution of files above indicates that the internship work was not limited to a single page or isolated feature branch. Instead, it covered high-traffic screens such as checkout, header navigation, product details, orders, and core admin pages that connect directly to operational behavior.",
    )

    add_subheading(doc, "High-Churn Files from Authored Work")
    for file_path, count in data["top_churn_files"]:
        add_bullet(doc, f"{file_path} ({count} total changed lines in authored history)")
    add_body(
        doc,
        "High churn in these files is expected because they are central to product experience or business workflows. Their presence in the report helps explain where the most iterative problem-solving occurred during the internship period.",
    )

    add_subheading(doc, "Representative Recent Commits")
    for commit in data["recent_commits"]:
        add_bullet(doc, commit)
    add_body(
        doc,
        "The recent commit sequence illustrates the kind of work I was actively handling toward the end of the review period: product presentation refinement, image/API integration, order flow correctness, responsive improvements, and linked fixes spanning multiple frontend modules.",
    )

    add_subheading(doc, "Case-Based Explanation of Work")
    for title, paragraph in data["case_studies"]:
        add_body(doc, f"{title}: {paragraph}")
        add_body(
            doc,
            f"The relevance of {title.lower()} to the project is that it combines interface concerns with functional correctness. In practice, such work required reading multiple related files, understanding the user journey end to end, and then making changes that preserved behavior across connected parts of the application.",
        )

    add_subheading(doc, "Challenges Faced During Internship")
    for challenge in data["challenges"]:
        add_bullet(doc, str(challenge))
    add_body(
        doc,
        "These challenges were valuable because they pushed me to work more like a real product engineer. I had to move beyond isolated coding tasks and instead think in terms of shared impact, flow continuity, user expectation, and system-level correctness.",
    )

    add_subheading(doc, "Professional Growth Through Contribution")
    add_body(
        doc,
        "The internship gradually improved my confidence in navigating existing code, tracing side effects of changes, and making practical decisions in a collaborative environment. This growth is visible in the spread of contribution across routes, components, state containers, and admin screens rather than a narrow concentration on only one feature area.",
    )
    add_body(
        doc,
        "Another major gain was learning to treat frontend quality as a combination of design clarity, logic correctness, and reliability under real user flows. This perspective is essential in ecommerce products, where the frontend is directly responsible for both trust and conversion.",
    )

    add_heading(doc, "5 LEARNING OUTCOME FROM 6 MONTHS INTERNSHIP")
    for point in data["learning_points"]:
        add_bullet(doc, str(point))
    add_body(
        doc,
        "One of the strongest learning outcomes was understanding how a real project evolves through repeated incremental improvements. Instead of writing code once and moving on, I experienced the cycle of initial implementation, bug fixing, refinement, regression awareness, and adaptation to changing requirements.",
    )
    add_body(
        doc,
        "I also learned the importance of structured debugging. Many issues that appeared visual at first were actually rooted in state handling, utility behavior, API data, or interactions between multiple reused modules. This reinforced the need to diagnose problems carefully before editing code.",
    )
    add_body(
        doc,
        "The internship improved my communication with the codebase itself. Reading commit history, tracing module relationships, and understanding how a shared component affects multiple screens made me far more comfortable working in an existing repository than building only from scratch.",
    )
    add_body(
        doc,
        "Finally, the experience strengthened my appreciation for delivery discipline. Correctness in pricing, responsiveness in product display, continuity in user journeys, and maintainability of shared components are all part of professional frontend engineering, and this internship provided practical exposure to each of those dimensions.",
    )
    add_subheading(doc, "Technical Learning")
    add_body(
        doc,
        "Technically, the internship made me more capable in React and Next.js development within an existing codebase. I improved in reading established code structures, tracing state interactions, diagnosing rendering issues, and making changes that stayed compatible with surrounding modules rather than breaking reuse.",
    )
    add_body(
        doc,
        "I also became more comfortable with frontend logic that is closely tied to business behavior. Cart totals, order calculations, membership display, wishlist state, and support-related flows showed me that frontend development often includes logic-sensitive responsibilities rather than only visual composition.",
    )

    add_subheading(doc, "Project and Product Learning")
    add_body(
        doc,
        "From a product perspective, I learned to think in terms of end-to-end user journeys. A good feature is not only one that renders correctly on a single page, but one that remains coherent from discovery to checkout to order follow-up. This mindset helped me evaluate my own work more carefully and understand why regressions in one area can damage the entire experience.",
    )
    add_body(
        doc,
        "Figure 4: Learning Outcome Framework. The internship strengthened my growth across four connected areas: technical implementation, debugging discipline, product understanding, and collaborative development practices.",
    )

    add_subheading(doc, "Professional Learning")
    add_body(
        doc,
        "Professionally, the internship improved my discipline in working within a collaborative repository. I learned to treat changes more carefully, respect existing structures, understand linked modules before editing them, and rely on history and evidence when describing my contribution. This is one of the most important shifts from academic coding to real project work.",
    )

    add_heading(doc, "6 CONCLUSION")
    add_body(
        doc,
        f"My internship work on {data['company_name']} gave me practical exposure to real frontend engineering in an active product codebase. I contributed across pages, components, state management, and flow correctness, especially in ecommerce journeys where user experience and pricing accuracy were critical.",
    )
    add_body(
        doc,
        "The experience improved my confidence in understanding existing systems, making changes safely in a collaborative repository, and translating functional requirements into maintainable frontend updates. It also strengthened my ability to connect design, implementation, debugging, and delivery in a professional development environment.",
    )
    add_body(
        doc,
        "The evidence gathered from repository history supports the conclusion that my contribution was broad, consistent, and frontend-centered. It covered customer-facing flows, shared UI modules, business-sensitive checkout logic, and supporting admin interfaces. This made the internship a meaningful bridge between academic learning and real software product work.",
    )
    add_body(
        doc,
        "Overall, the internship was valuable not only because of the amount of work completed, but also because of the kind of work completed. It required responsibility, attention to detail, and growing independence in a production-style engineering environment, which together made it a strong practical learning experience.",
    )

    add_heading(doc, "REFERENCES")
    add_bullet(doc, "Project repository history and authored Git commits from the BogEcom workspace.")
    add_bullet(doc, "Project source modules under frontend/client, frontend/admin, and server.")
    add_bullet(doc, "Next.js and React documentation used as standard technology references during development.")
    add_bullet(doc, "Internal file and commit evidence summarized from version-control records in the local workspace.")

    section = doc.sections[-1]
    footer = section.footer.paragraphs[0]
    footer.alignment = WD_ALIGN_PARAGRAPH.CENTER
    add_page_number(footer)

    try:
        doc.save(OUTPUT_DOCX)
        print(f"Generated: {OUTPUT_DOCX}")
    except PermissionError:
        doc.save(FALLBACK_OUTPUT_DOCX)
        print(f"Generated: {FALLBACK_OUTPUT_DOCX}")


def clear_slide(slide) -> None:
    shape_tree = slide.shapes._spTree
    for shape in list(slide.shapes):
        if not shape.is_placeholder:
            shape_tree.remove(shape._element)


def set_textbox(slide, left, top, width, height, text, size=24, bold=False, color=(20, 40, 80), align=PP_ALIGN.LEFT):
    textbox = slide.shapes.add_textbox(left, top, width, height)
    frame = textbox.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    p = frame.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0]
    run.font.name = "Calibri"
    run.font.size = PPTPt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*color)
    return textbox


def add_bullets_to_frame(frame, items, size=22, color=(50, 50, 50)):
    frame.clear()
    for idx, item in enumerate(items):
        p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        p.text = item
        p.level = 0
        p.bullet = True
        p.alignment = PP_ALIGN.LEFT
        for run in p.runs:
            run.font.name = "Calibri"
            run.font.size = PPTPt(size)
            run.font.color.rgb = RGBColor(*color)


def add_title_and_body_slide(prs: Presentation, title: str, bullets: list[str]):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    title_p = slide.shapes.title.text_frame.paragraphs[0]
    for run in title_p.runs:
        run.font.bold = True
        run.font.size = PPTPt(28)
        run.font.color.rgb = RGBColor(20, 40, 80)
    body = slide.placeholders[1].text_frame
    add_bullets_to_frame(body, bullets)
    return slide


def add_highlight_boxes(slide, left, top, items):
    box_width = PPTInches(2.1)
    gap = PPTInches(0.2)
    for idx, (value, label) in enumerate(items):
        x = left + idx * (box_width + gap)
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, top, box_width, PPTInches(1.2))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(232, 239, 255)
        shape.line.color.rgb = RGBColor(120, 150, 210)
        tf = shape.text_frame
        tf.clear()
        p1 = tf.paragraphs[0]
        p1.text = value
        p1.alignment = PP_ALIGN.CENTER
        r1 = p1.runs[0]
        r1.font.bold = True
        r1.font.size = PPTPt(24)
        r1.font.color.rgb = RGBColor(20, 40, 80)
        p2 = tf.add_paragraph()
        p2.text = label
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.runs[0]
        r2.font.size = PPTPt(13)
        r2.font.color.rgb = RGBColor(70, 70, 70)


def build_presentation(data: dict[str, object]) -> None:
    prs = Presentation(str(TEMPLATE_PPT))

    cover = prs.slides[0]
    for shape in cover.shapes:
        text = getattr(shape, "text", "")
        if "Mid Review Internship Presentation" in text:
            shape.text = "Mid Review Internship Presentation"
        elif "Presented by" in text:
            shape.text = (
                "Presented by\n"
                f"{data['student_name']}\n"
                f"Role: {data['role']}\n"
                f"Project: {data['company_name']}"
            )
        elif "Faculty Internship guide" in text:
            shape.text = f"Faculty Internship guide:-\n{data['faculty_guide']}"
        elif "Industrial guide" in text:
            shape.text = f"Industrial guide:-\n{data['industry_guide']}"
        elif "2025-2026" in text:
            shape.text = "2025-2026"
    set_textbox(
        cover,
        PPTInches(0.9),
        PPTInches(2.2),
        PPTInches(3.6),
        PPTInches(2.0),
        f"Presented by\n{data['student_name']}\n{data['role']}\n{data['company_name']}",
        20,
        True,
        (20, 40, 80),
        PP_ALIGN.LEFT,
    )
    set_textbox(
        cover,
        PPTInches(0.9),
        PPTInches(4.5),
        PPTInches(3.2),
        PPTInches(0.5),
        f"Reg. No.: {data['registration_number']}",
        16,
        False,
        (70, 70, 70),
        PP_ALIGN.LEFT,
    )
    if LOGO_PATH.exists():
        cover.shapes.add_picture(str(LOGO_PATH), PPTInches(8.3), PPTInches(1.3), height=PPTInches(1.2))

    contents = prs.slides[1]
    for shape in contents.shapes:
        if getattr(shape, "text", "").strip() == "Contents":
            continue
        if getattr(shape, "text", "").strip():
            shape.text = (
                "Company / Project Profile\n\n"
                "Internship Role\n\n"
                "Project Introduction\n\n"
                "Work Evidence and Metrics\n\n"
                "Roles and Responsibilities\n\n"
                "Tools Used\n\n"
                "Methodology\n\n"
                "Outcomes and Learning\n\n"
                "Conclusion"
            )

    add_title_and_body_slide(
        prs,
        "Company / Project Profile",
        [
            f"Project: {data['company_name']}",
            "Type: Full-stack ecommerce platform with client and admin interfaces",
            "Core areas: product discovery, cart, checkout, orders, membership, support, and admin operations",
            "My primary area of contribution: frontend implementation and flow stabilization",
        ],
    )

    add_title_and_body_slide(
        prs,
        "Internship Role",
        [
            f"Role executed: {data['role']}",
            "Worked mainly on frontend/client and selected frontend/admin modules",
            "Focused on production-style bug fixing, UI refinement, and API-integrated features",
            f"Active authored frontend work window: {data['start_date']} to {data['end_date']}",
        ],
    )

    add_title_and_body_slide(
        prs,
        "Project Introduction",
        [
            "The project uses React and Next.js to deliver a modern ecommerce experience.",
            "The frontend includes reusable components, shared context state, and API-connected pages.",
            "My work concentrated on solving practical issues in user-facing journeys and keeping the UI consistent.",
        ],
    )

    metrics_slide = prs.slides.add_slide(prs.slide_layouts[5])
    set_textbox(metrics_slide, PPTInches(0.5), PPTInches(0.2), PPTInches(9), PPTInches(0.5), "Work Evidence and Metrics", 28, True)
    add_highlight_boxes(
        metrics_slide,
        PPTInches(0.5),
        PPTInches(1.1),
        [
            (str(data["frontend_commit_count"]), "Frontend-related commits"),
            (str(data["unique_files"]), "Unique files touched"),
            (str(data["overall_commit_count"]), "Overall repo commits by my identities"),
            ("2", "Frontend apps contributed"),
        ],
    )
    set_textbox(metrics_slide, PPTInches(0.6), PPTInches(2.7), PPTInches(4.3), PPTInches(0.4), "Top contribution areas", 20, True)
    set_textbox(
        metrics_slide,
        PPTInches(0.6),
        PPTInches(3.1),
        PPTInches(4.3),
        PPTInches(2.1),
        "\n".join([f"{label}: {count} files" for label, count in data["category_counter"].most_common(5)]),
        17,
        False,
    )
    set_textbox(metrics_slide, PPTInches(5.1), PPTInches(2.7), PPTInches(4.0), PPTInches(0.4), "Recent commit examples", 20, True)
    set_textbox(
        metrics_slide,
        PPTInches(5.1),
        PPTInches(3.1),
        PPTInches(4.0),
        PPTInches(2.6),
        "\n".join(data["recent_commits"][:5]),
        14,
        False,
    )

    add_title_and_body_slide(prs, "Roles and Responsibilities", list(data["work_summary"]))
    add_title_and_body_slide(prs, "Tools Used", list(data["tools"]))
    add_title_and_body_slide(prs, "Methodology", list(data["methodology"]))
    add_title_and_body_slide(prs, "Outcomes and Learning", list(data["outcomes"]) + list(data["learning_points"][:2]))
    add_title_and_body_slide(
        prs,
        "Conclusion",
        [
            "The internship gave me direct experience with a real ecommerce codebase and team workflow.",
            "My major contribution area was frontend development, especially UI behavior, page flows, and pricing-sensitive journeys.",
            "The work improved my ability to debug, implement, and ship changes in an existing production-style repository.",
        ],
    )

    try:
        prs.save(str(OUTPUT_PPT))
        print(f"Generated: {OUTPUT_PPT}")
    except PermissionError:
        prs.save(str(FALLBACK_OUTPUT_PPT))
        print(f"Generated: {FALLBACK_OUTPUT_PPT}")


if __name__ == "__main__":
    stats = get_repo_stats()
    build_report(stats)
    build_presentation(stats)
