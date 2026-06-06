from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt


TEMPLATE = Path(r"C:\Users\piyush songara\Downloads\PPT Format.pptx")
OUTPUT = Path(r"C:\Users\piyush songara\Downloads\Netparam_Mid_Review_Presentation.pptx")

BG = Path(r"D:\AnanyaBoutique\output\ppt-template-media\image1.jpeg")
CAMPUS = Path(r"D:\AnanyaBoutique\output\ppt-template-media\image2.png")
JECRC_LOGO = Path(r"D:\AnanyaBoutique\output\ppt-template-media\image3.jpeg")
COMPANY_LOGO = Path(r"D:\AnanyaBoutique\output\friend-report-media\netparam_logo_generated.png")
ARCH = Path(r"D:\AnanyaBoutique\output\netparam-ppt-media\image3.png")
FLOW = Path(r"D:\AnanyaBoutique\output\netparam-ppt-media\image4.png")
TEST = Path(r"D:\AnanyaBoutique\output\netparam-ppt-media\image5.png")
LEARN = Path(r"D:\AnanyaBoutique\output\netparam-ppt-media\image6.png")

RED = RGBColor(230, 24, 34)
BLACK = RGBColor(0, 0, 0)


def set_text(tf, text, size=20, bold=False, color=BLACK, align=PP_ALIGN.LEFT, font_name="Times New Roman"):
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font_name
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return p


def add_bullets(tf, title, bullets, title_size=24, body_size=18):
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = title
    r.font.name = "Times New Roman"
    r.font.size = Pt(title_size)
    r.font.bold = True
    r.font.color.rgb = BLACK
    for item in bullets:
        p = tf.add_paragraph()
        p.text = item
        p.level = 0
        p.alignment = PP_ALIGN.LEFT
        for r in p.runs:
            r.font.name = "Times New Roman"
            r.font.size = Pt(body_size)
            r.font.color.rgb = BLACK


def add_bg(slide):
    slide.shapes.add_picture(str(BG), 0, 0, width=prs.slide_width, height=prs.slide_height)


def add_title(slide, title):
    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.32), Inches(8.2), Inches(0.55))
    set_text(box.text_frame, title, size=22, bold=True, color=RGBColor(255, 255, 255), align=PP_ALIGN.CENTER)
    return box


def add_textbox(slide, left, top, width, height, text, size=18, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    set_text(box.text_frame, text, size=size, bold=bold, align=align)
    return box


def add_body_box(slide, left, top, width, height, title, bullets):
    rect = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(255, 255, 255)
    rect.line.color.rgb = RGBColor(215, 215, 215)
    add_bullets(rect.text_frame, title, bullets)
    return rect


prs = Presentation(str(TEMPLATE))
slide_w = prs.slide_width
slide_h = prs.slide_height

# Slide 1: cover
s1 = prs.slides[0]
for shape in list(s1.shapes):
    if hasattr(shape, "text_frame"):
        txt = shape.text.strip()
        if "Mid Review Internship Presentation" in txt:
            set_text(shape.text_frame, "Mid Review Internship Presentation", size=23, color=RGBColor(128, 128, 128), align=PP_ALIGN.CENTER)
        elif "Presented by" in txt:
            text = "Presented by\n[Your Name]\nUnder the guidance of"
            set_text(shape.text_frame, text, size=18, bold=False, align=PP_ALIGN.CENTER)
            for p in shape.text_frame.paragraphs:
                for r in p.runs:
                    if "[Your Name]" in r.text:
                        r.font.bold = True
                        r.font.size = Pt(20)
        elif "Faculty Internship guide" in txt:
            set_text(shape.text_frame, "Faculty Internship guide:-\n[Faculty Guide Name]", size=16, bold=False, align=PP_ALIGN.CENTER)
        elif "<COMPANY LOGO>" in txt:
            set_text(shape.text_frame, "Industrial guide:-\n[Industry Guide Name]\n[Industry Guide Designation]", size=16, align=PP_ALIGN.CENTER)
        elif "2025-2026" in txt:
            set_text(shape.text_frame, "2025-2026", size=18, bold=False, align=PP_ALIGN.CENTER)
        elif "DEPARTMENT OF" in txt:
            set_text(shape.text_frame, "DEPARTMENT OF\nComputer Science", size=22, bold=True, color=RGBColor(255, 255, 255), align=PP_ALIGN.CENTER)

# Replace company logo placeholder area
s1.shapes.add_picture(str(COMPANY_LOGO), Inches(9.2), Inches(3.1), width=Inches(1.6), height=Inches(1.45))

# Slide 2: contents
s2 = prs.slides[1]
for shape in s2.shapes:
    if hasattr(shape, "text_frame"):
        txt = shape.text.strip()
        if txt == "Contents":
            set_text(shape.text_frame, "Contents", size=24, bold=True, color=RGBColor(255, 255, 255), align=PP_ALIGN.CENTER)
        elif "Company Profile" in txt:
            contents = (
                "Company Profile\n\n"
                "Internship Role\n\n"
                "Project Introduction\n\n"
                "Tools Used\n\n"
                "Methodology\n\n"
                "Testing and Validation\n\n"
                "Outcomes and Learning\n\n"
                "Conclusion"
            )
            set_text(shape.text_frame, contents, size=20, align=PP_ALIGN.LEFT)


def new_content_slide(title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, title)
    return slide


# Slide 3 company profile
slide = new_content_slide("Company Profile")
add_body_box(
    slide,
    Inches(0.65),
    Inches(1.25),
    Inches(5.2),
    Inches(3.0),
    "About Netparam Technologies",
    [
        "Internship organization for the reported project work.",
        "Provided a professional environment for applied machine learning development.",
        "Supported implementation, documentation, and review preparation in the ALCHEMIST workspace.",
    ],
)
slide.shapes.add_picture(str(COMPANY_LOGO), Inches(6.5), Inches(1.55), width=Inches(2.3), height=Inches(2.1))
add_body_box(
    slide,
    Inches(5.95),
    Inches(3.15),
    Inches(3.65),
    Inches(1.35),
    "Internship Role",
    [
        "Worked on a machine learning classification system with model training, testing, and Flask integration.",
    ],
)

# Slide 4 project introduction
slide = new_content_slide("Project Introduction")
add_body_box(
    slide,
    Inches(0.55),
    Inches(1.15),
    Inches(4.25),
    Inches(4.2),
    "Perfume Notes Classification System",
    [
        "Predicts fragrance families from perfume note descriptions.",
        "Uses multi-label output: Floral, Woody, Citrus, and Oriental.",
        "Converts text notes into features and serves predictions through Flask.",
        "Designed as a compact end-to-end internship project.",
    ],
)
slide.shapes.add_picture(str(ARCH), Inches(4.95), Inches(1.2), width=Inches(4.4), height=Inches(4.1))

# Slide 5 tools
slide = new_content_slide("Tools Used")
add_body_box(
    slide,
    Inches(0.7),
    Inches(1.35),
    Inches(8.6),
    Inches(3.7),
    "Technology Stack",
    [
        "Python: primary language for model training and application logic",
        "Pandas and CSV handling: dataset preparation and structuring",
        "scikit-learn: TF-IDF, train-test split, One-vs-Rest, Naive Bayes",
        "Flask: lightweight web application and prediction interface",
        "Pickle: model and vectorizer persistence",
        "HTML, CSS, JavaScript: user-facing prediction page",
    ],
)

# Slide 6 methodology
slide = new_content_slide("Methodology")
slide.shapes.add_picture(str(FLOW), Inches(0.85), Inches(1.18), width=Inches(8.35), height=Inches(4.65))

# Slide 7 testing
slide = new_content_slide("Testing and Validation")
slide.shapes.add_picture(str(TEST), Inches(0.85), Inches(1.18), width=Inches(8.35), height=Inches(4.65))

# Slide 8 outcomes
slide = new_content_slide("Outcomes and Learning")
slide.shapes.add_picture(str(LEARN), Inches(0.65), Inches(1.15), width=Inches(4.7), height=Inches(4.55))
add_body_box(
    slide,
    Inches(5.5),
    Inches(1.4),
    Inches(3.55),
    Inches(3.7),
    "Key Outcomes",
    [
        "Improved understanding of machine learning fundamentals.",
        "Learned text preprocessing and feature engineering.",
        "Built confidence in Flask-based model serving.",
        "Developed documentation and presentation discipline.",
        "Understood end-to-end project workflow from data to interface.",
    ],
)

# Slide 9 roles and responsibilities
slide = new_content_slide("Roles and Responsibilities")
add_body_box(
    slide,
    Inches(0.65),
    Inches(1.3),
    Inches(4.3),
    Inches(3.95),
    "Work Done",
    [
        "Selected project scope and defined fragrance categories.",
        "Prepared and labeled the perfume notes dataset.",
        "Implemented TF-IDF vectorization and classifier training.",
        "Saved trained artifacts for reuse.",
        "Integrated prediction flow into a Flask web interface.",
    ],
)
add_body_box(
    slide,
    Inches(5.15),
    Inches(1.3),
    Inches(4.15),
    Inches(3.95),
    "Project Evidence",
    [
        "Repository contains train.py, app.py, perfumes.csv, and model files.",
        "Manual testing performed on sample note combinations.",
        "System designed for local demonstration and future deployment.",
        "Report prepared in prescribed university format.",
    ],
)

# Slide 10 conclusion
slide = new_content_slide("Conclusion")
add_body_box(
    slide,
    Inches(0.85),
    Inches(1.45),
    Inches(8.0),
    Inches(2.95),
    "Conclusion",
    [
        "The Perfume Notes Classification System demonstrates a complete workflow from dataset design to user-facing prediction.",
        "The internship strengthened practical skills in machine learning, text processing, model persistence, and Flask integration.",
        "The project is compact, reproducible, and suitable for academic review and technical demonstration.",
    ],
)
add_textbox(slide, Inches(2.7), Inches(4.8), Inches(4.3), Inches(0.5), "Thank You", size=24, bold=True, align=PP_ALIGN.CENTER)

prs.save(str(OUTPUT))
print(OUTPUT)
