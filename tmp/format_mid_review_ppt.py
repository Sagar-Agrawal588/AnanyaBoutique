from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Pt


PPT_PATH = Path(r"d:\BogEcom\Piyush_Songara_Mid_Review_Presentation.pptx")


def normalize_text(shape) -> None:
    text = getattr(shape, "text", None)
    if not text:
        return
    cleaned = text.replace("\x0b", "\n").replace("\r\n", "\n").strip()
    if cleaned != text.strip():
        shape.text = cleaned


def style_paragraph(paragraph, *, font_size: int, align, bold: bool = False) -> None:
    paragraph.alignment = align
    for run in paragraph.runs:
        run.font.name = "Times New Roman"
        run.font.size = Pt(font_size)
        run.font.bold = bold


def style_text_frame(shape, slide_index: int) -> None:
    if not hasattr(shape, "text_frame"):
        return

    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Pt(6)
    tf.margin_right = Pt(6)
    tf.margin_top = Pt(4)
    tf.margin_bottom = Pt(4)

    text = shape.text.strip()
    if not text:
        return

    if slide_index == 1:
        if "Mid Review Internship Presentation" in text:
            for p in tf.paragraphs:
                style_paragraph(p, font_size=24, align=PP_ALIGN.CENTER, bold=True)
        elif "Presented by" in text or "Reg. No.:" in text:
            for p in tf.paragraphs:
                style_paragraph(p, font_size=18, align=PP_ALIGN.LEFT, bold=False)
        elif "Faculty Internship guide" in text or "Industrial guide" in text:
            for p in tf.paragraphs:
                style_paragraph(p, font_size=16, align=PP_ALIGN.LEFT, bold=False)
        else:
            for p in tf.paragraphs:
                style_paragraph(p, font_size=16, align=PP_ALIGN.CENTER, bold=False)
        return

    is_title = text in {
        "Contents",
        "Company / Project Profile",
        "Internship Role",
        "Project Introduction",
        "Work Evidence and Metrics",
        "Roles and Responsibilities",
        "Tools Used",
        "Methodology",
        "Outcomes and Learning",
        "Conclusion",
    }

    if is_title:
        for p in tf.paragraphs:
            style_paragraph(p, font_size=24, align=PP_ALIGN.CENTER, bold=True)
        return

    if slide_index == 6 and (
        "Frontend-related commits" in text
        or "Unique files touched" in text
        or "Overall repo commits by my identities" in text
        or "Frontend apps contributed" in text
    ):
        first = True
        for p in tf.paragraphs:
            if not p.text.strip():
                continue
            if first:
                style_paragraph(p, font_size=22, align=PP_ALIGN.CENTER, bold=True)
                first = False
            else:
                style_paragraph(p, font_size=12, align=PP_ALIGN.CENTER, bold=False)
        return

    if slide_index == 6 and (
        "Server and API:" in text
        or "Pages and flows:" in text
        or text.startswith("e0c4845")
        or text.startswith("9c61a22")
    ):
        for p in tf.paragraphs:
            if p.text.strip():
                style_paragraph(p, font_size=12, align=PP_ALIGN.LEFT, bold=False)
        return

    # Metric cards on slide 6
    if (
        slide_index == 6
        and "\n" in text
        and len(tf.paragraphs) >= 2
        and "files" not in text
        and "e0c4845" not in text
        and "Recent commit examples" not in text
    ):
        first = True
        for p in tf.paragraphs:
            if not p.text.strip():
                continue
            if first:
                style_paragraph(p, font_size=22, align=PP_ALIGN.CENTER, bold=True)
                first = False
            else:
                style_paragraph(p, font_size=12, align=PP_ALIGN.CENTER, bold=False)
        return

    if slide_index == 2:
        for p in tf.paragraphs:
            if p.text.strip():
                style_paragraph(p, font_size=18, align=PP_ALIGN.LEFT, bold=False)
        return

    if slide_index == 6 and (
        text.startswith("Top contribution areas")
        or text.startswith("Recent commit examples")
        or text.startswith("Server and API:")
        or text.startswith("Pages and flows:")
        or text.startswith("UI components:")
        or text.startswith("Config and assets:")
        or text.startswith("Utilities and API:")
        or text.startswith("State management:")
        or text.startswith("e0c4845")
    ):
        for idx, p in enumerate(tf.paragraphs):
            if not p.text.strip():
                continue
            style_paragraph(
                p,
                font_size=16 if idx == 0 and ("Top contribution areas" in text or "Recent commit examples" in text) else 12,
                align=PP_ALIGN.LEFT,
                bold=idx == 0 and ("Top contribution areas" in text or "Recent commit examples" in text),
            )
        return

    for p in tf.paragraphs:
        if p.text.strip():
            style_paragraph(p, font_size=18, align=PP_ALIGN.LEFT, bold=False)


def main() -> None:
    prs = Presentation(str(PPT_PATH))
    for slide_index, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            normalize_text(shape)
            style_text_frame(shape, slide_index)
    prs.save(str(PPT_PATH))
    print(PPT_PATH)


if __name__ == "__main__":
    main()
